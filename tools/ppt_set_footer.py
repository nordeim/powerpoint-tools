#!/usr/bin/env python3
"""
PowerPoint Set Footer Tool v3.1.0
Configure slide footer with Dual Strategy (Placeholder + Text Box Fallback).

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.0

Usage:
    uv run tools/ppt_set_footer.py --file deck.pptx --text "Company © 2024" --json
    uv run tools/ppt_set_footer.py --file deck.pptx --text "Confidential" --show-number --json

Exit Codes:
    0: Success
    1: Error occurred
"""

import sys
import os

sys.stderr = open(os.devnull, 'w')

import json
import argparse
from pathlib import Path
from typing import Dict, Any, Set

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import PowerPointAgent

try:
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError:
    class PP_PLACEHOLDER:
        FOOTER = 15
        SLIDE_NUMBER = 13

__version__ = "3.1.0"


def set_footer(
    filepath: Path,
    text: str = None,
    show_number: bool = False
) -> Dict[str, Any]:
    """
    Set footer on slides using Dual Strategy.
    
    Args:
        filepath: Path to PowerPoint file (.pptx)
        text: Footer text
        show_number: Whether to show slide numbers
        
    Returns:
        Dict with results
        
    Raises:
        FileNotFoundError: If file doesn't exist
        ValueError: If file format invalid
    """
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    
    if filepath.suffix.lower() != '.pptx':
        raise ValueError("Only .pptx files are supported")
    
    slide_indices_updated: Set[int] = set()
    method_used = None
    
    with PowerPointAgent(filepath) as agent:
        agent.open(filepath)
        
        version_before = agent.get_presentation_version()
        
        # Strategy 1: Try placeholders on slide masters
        try:
            for master in agent.prs.slide_masters:
                for layout in master.slide_layouts:
                    for shape in layout.placeholders:
                        try:
                            if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                                if text:
                                    shape.text = text
                        except Exception:
                            pass
        except Exception:
            pass
        
        # Try placeholders on slides
        for slide_idx, slide in enumerate(agent.prs.slides):
            try:
                for shape in slide.placeholders:
                    try:
                        if shape.placeholder_format.type == PP_PLACEHOLDER.FOOTER:
                            if text:
                                shape.text = text
                            slide_indices_updated.add(slide_idx)
                    except Exception:
                        pass
            except Exception:
                pass
        
        # Strategy 2: Fallback to text boxes if placeholders didn't work
        if len(slide_indices_updated) == 0:
            method_used = "text_box"
            for slide_idx in range(len(agent.prs.slides)):
                try:
                    if text:
                        agent.add_text_box(
                            slide_index=slide_idx,
                            text=text,
                            position={"left": "5%", "top": "92%"},
                            size={"width": "60%", "height": "5%"},
                            font_size=10,
                            color="#595959"
                        )
                        slide_indices_updated.add(slide_idx)
                    if show_number:
                        agent.add_text_box(
                            slide_index=slide_idx,
                            text=str(slide_idx + 1),
                            position={"left": "92%", "top": "92%"},
                            size={"width": "5%", "height": "5%"},
                            font_size=10,
                            color="#595959"
                        )
                        slide_indices_updated.add(slide_idx)
                except Exception:
                    pass
        else:
            method_used = "placeholder"
        
        agent.save()
        
        version_after = agent.get_presentation_version()
    
    return {
        "status": "success" if len(slide_indices_updated) > 0 else "warning",
        "file": str(filepath.resolve()),
        "method_used": method_used,
        "slides_updated": len(slide_indices_updated),
        "footer_text": text,
        "show_number": show_number,
        "presentation_version_before": version_before,
        "presentation_version_after": version_after,
        "tool_version": __version__
    }


def main():
    parser = argparse.ArgumentParser(
        description="Set slide footer with text and/or page numbers",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Set footer text
  uv run tools/ppt_set_footer.py --file deck.pptx --text "Company © 2024" --json

  # Add page numbers
  uv run tools/ppt_set_footer.py --file deck.pptx --show-number --json

  # Both footer text and page numbers
  uv run tools/ppt_set_footer.py --file deck.pptx --text "Confidential" --show-number --json

Strategy:
  1. Tries to use slide placeholders first (preserves template formatting)
  2. Falls back to text boxes if placeholders not available
        """
    )
    
    parser.add_argument('--file', required=True, type=Path, help='PowerPoint file path (.pptx)')
    parser.add_argument('--text', help='Footer text')
    parser.add_argument('--show-number', action='store_true', help='Show slide numbers')
    parser.add_argument('--show-date', action='store_true', help='Show date (placeholder only)')
    parser.add_argument('--json', action='store_true', default=True, help='Output JSON (default: true)')
    
    args = parser.parse_args()
    
    try:
        result = set_footer(
            filepath=args.file,
            text=args.text,
            show_number=args.show_number
        )
        print(json.dumps(result, indent=2))
        sys.exit(0)
        
    except FileNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "FileNotFoundError",
            "suggestion": "Verify file path exists and is accessible."
        }
        print(json.dumps(error_result, indent=2))
        sys.exit(1)
        
    except ValueError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ValueError",
            "suggestion": "Ensure file has .pptx extension."
        }
        print(json.dumps(error_result, indent=2))
        sys.exit(1)
        
    except Exception as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "tool_version": __version__
        }
        print(json.dumps(error_result, indent=2))
        sys.exit(1)


if __name__ == "__main__":
    main()
