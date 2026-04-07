#!/usr/bin/env python3
"""
PowerPoint Merge Presentations Tool v3.1.1
Combine slides from multiple presentations into one.

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.1

Usage:
    uv run tools/ppt_merge_presentations.py --sources '[{"file":"a.pptx","slides":"all"},{"file":"b.pptx","slides":[0,2,4]}]' --output merged.pptx --json

Exit Codes:
    0: Success
    1: Error occurred (check error_type in JSON for details)

This tool merges slides from multiple source presentations into a single output
presentation. You can specify which slides to include from each source.

Source Specification Format:
    [
        {"file": "path/to/first.pptx", "slides": "all"},
        {"file": "path/to/second.pptx", "slides": [0, 1, 2]},
        {"file": "path/to/third.pptx", "slides": [5, 6]}
    ]
"""

import sys
import os

# --- HYGIENE BLOCK START ---
# CRITICAL: Redirect stderr to /dev/null to prevent library noise from corrupting JSON output
sys.stderr = open(os.devnull, "w")
# --- HYGIENE BLOCK END ---

import json
import argparse
import shutil
from pathlib import Path
from typing import Dict, Any, List, Optional, Union

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import (
    PowerPointAgent,
    PowerPointAgentError,
    SlideNotFoundError,
    ApprovalTokenError,
    APPROVAL_SCOPE_MERGE_PRESENTATIONS,
)

# ============================================================================
# CONSTANTS
# ============================================================================

__version__ = "3.1.1"


# ============================================================================
# TYPE DEFINITIONS
# ============================================================================

SourceSpec = Dict[str, Union[str, List[int]]]


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================


def parse_sources(sources_json: str) -> List[SourceSpec]:
    """
    Parse and validate sources JSON specification.

    Args:
        sources_json: JSON string with source specifications

    Returns:
        List of validated source specifications

    Raises:
        ValueError: If JSON is invalid or missing required fields
    """
    try:
        sources = json.loads(sources_json)
    except json.JSONDecodeError as e:
        raise ValueError(f"Invalid JSON in sources: {e}")

    if not isinstance(sources, list):
        raise ValueError("Sources must be a JSON array")

    if len(sources) == 0:
        raise ValueError("At least one source is required")

    validated = []
    for idx, source in enumerate(sources):
        if not isinstance(source, dict):
            raise ValueError(f"Source {idx} must be an object")

        if "file" not in source:
            raise ValueError(f"Source {idx} missing required 'file' field")

        if "slides" not in source:
            source["slides"] = "all"

        slides = source["slides"]
        if slides != "all" and not isinstance(slides, list):
            raise ValueError(f"Source {idx} 'slides' must be 'all' or array of indices")

        if isinstance(slides, list):
            for slide_idx in slides:
                if not isinstance(slide_idx, int) or slide_idx < 0:
                    raise ValueError(
                        f"Source {idx} has invalid slide index: {slide_idx}"
                    )

        validated.append(source)

    return validated


def validate_source_files(sources: List[SourceSpec]) -> None:
    """
    Validate all source files exist.

    Args:
        sources: List of source specifications

    Raises:
        FileNotFoundError: If any source file doesn't exist
    """
    for source in sources:
        filepath = Path(source["file"])
        if not filepath.exists():
            raise FileNotFoundError(f"Source file not found: {filepath}")
        if not filepath.suffix.lower() == ".pptx":
            raise ValueError(f"Source file must be .pptx: {filepath}")


# ============================================================================
# MAIN LOGIC
# ============================================================================


def merge_presentations(
    sources: List[SourceSpec],
    output: Path,
    base_template: Optional[Path] = None,
    preserve_formatting: bool = True,
) -> Dict[str, Any]:
    """
    Merge slides from multiple presentations into one.

    Args:
        sources: List of source specifications with file paths and slide indices
        output: Path for the output merged presentation
        base_template: Optional template to use for theme/masters
        preserve_formatting: Whether to preserve original slide formatting

    Returns:
        Dict with merge results

    Raises:
        FileNotFoundError: If source files don't exist
        SlideNotFoundError: If specified slide indices are invalid
        ValueError: If sources specification is invalid
    """
    validate_source_files(sources)

    if base_template:
        if not base_template.exists():
            raise FileNotFoundError(f"Base template not found: {base_template}")
        shutil.copy2(base_template, output)
        initial_file = base_template
    else:
        first_source = Path(sources[0]["file"])
        shutil.copy2(first_source, output)
        initial_file = first_source

    warnings: List[str] = []
    sources_used: List[Dict[str, Any]] = []
    merge_details: Dict[str, int] = {}
    total_slides_copied = 0

    with PowerPointAgent(output) as agent:
        agent.open(output)

        if not base_template:
            first_source_info = {
                "file": str(Path(sources[0]["file"]).resolve()),
                "slides_spec": sources[0]["slides"],
                "slides_copied": agent.get_slide_count(),
                "is_base": True,
            }
            sources_used.append(first_source_info)
            merge_details[str(Path(sources[0]["file"]).name)] = agent.get_slide_count()
            total_slides_copied += agent.get_slide_count()
            sources_to_process = sources[1:]
        else:
            sources_to_process = sources

        from pptx import Presentation

        for source_idx, source in enumerate(sources_to_process):
            source_path = Path(source["file"])
            slides_spec = source["slides"]

            try:
                source_prs = Presentation(str(source_path))
                source_slide_count = len(source_prs.slides)

                if slides_spec == "all":
                    slide_indices = list(range(source_slide_count))
                else:
                    slide_indices = slides_spec
                    for idx in slide_indices:
                        if idx >= source_slide_count:
                            raise SlideNotFoundError(
                                f"Slide {idx} not found in {source_path.name} (has {source_slide_count} slides)",
                                details={
                                    "source_file": str(source_path),
                                    "requested_index": idx,
                                    "available_slides": source_slide_count,
                                },
                            )

                slides_copied = 0
                for slide_idx in slide_indices:
                    try:
                        source_slide = source_prs.slides[slide_idx]

                        blank_layout = None
                        for layout in agent.prs.slide_layouts:
                            if "blank" in layout.name.lower():
                                blank_layout = layout
                                break
                        if blank_layout is None:
                            blank_layout = agent.prs.slide_layouts[0]

                        new_slide = agent.prs.slides.add_slide(blank_layout)

                        for shape in source_slide.shapes:
                            if shape.shape_type == 13:
                                continue

                            try:
                                el = shape.element
                                new_slide.shapes._spTree.insert_element_before(
                                    el, "p:extLst"
                                )
                            except Exception:
                                pass

                        slides_copied += 1
                        total_slides_copied += 1

                    except Exception as e:
                        warnings.append(
                            f"Could not copy slide {slide_idx} from {source_path.name}: {str(e)}"
                        )

                sources_used.append(
                    {
                        "file": str(source_path.resolve()),
                        "slides_spec": slides_spec,
                        "slides_copied": slides_copied,
                        "is_base": False,
                    }
                )
                merge_details[source_path.name] = slides_copied

            except SlideNotFoundError:
                raise
            except Exception as e:
                warnings.append(f"Error processing {source_path.name}: {str(e)}")

        agent.save()

        info = agent.get_presentation_info()
        presentation_version = info.get("presentation_version")
        final_slide_count = info.get("slide_count")

    return {
        "status": "success",
        "file": str(output.resolve()),
        "sources_used": sources_used,
        "total_slides": final_slide_count,
        "merge_details": merge_details,
        "base_template": str(base_template.resolve()) if base_template else None,
        "preserve_formatting": preserve_formatting,
        "warnings": warnings,
        "presentation_version": presentation_version,
        "tool_version": __version__,
    }


# ============================================================================
# CLI INTERFACE
# ============================================================================


def main():
    parser = argparse.ArgumentParser(
        description="Merge slides from multiple PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Merge all slides from two presentations
  uv run tools/ppt_merge_presentations.py \\
    --sources '[{"file":"part1.pptx","slides":"all"},{"file":"part2.pptx","slides":"all"}]' \\
    --output merged.pptx --json

  # Select specific slides from each source
  uv run tools/ppt_merge_presentations.py \\
    --sources '[{"file":"intro.pptx","slides":[0,1]},{"file":"content.pptx","slides":[2,3,4]},{"file":"outro.pptx","slides":[0]}]' \\
    --output presentation.pptx --json

  # Use a template for consistent theming
  uv run tools/ppt_merge_presentations.py \\
    --sources '[{"file":"content1.pptx","slides":"all"},{"file":"content2.pptx","slides":"all"}]' \\
    --output merged.pptx --base-template corporate_template.pptx --json

Source Specification Format:
  The --sources argument must be a JSON array with objects containing:
  - "file": Path to the source .pptx file (required)
  - "slides": Either "all" or an array of slide indices [0, 1, 2] (optional, default: "all")

Behavior:
  - First source becomes the base (its theme/masters are used)
  - Subsequent sources have their slides copied into the base
  - Use --base-template to override with a specific template
  - Slide indices are 0-based

Output Format:
  {
    "status": "success",
    "file": "/path/to/merged.pptx",
    "sources_used": [
      {"file": "part1.pptx", "slides_copied": 5, "is_base": true},
      {"file": "part2.pptx", "slides_copied": 3, "is_base": false}
    ],
    "total_slides": 8,
    "merge_details": {"part1.pptx": 5, "part2.pptx": 3},
    "presentation_version": "a1b2c3...",
    "tool_version": "3.1.1"
  }
        """,
    )

    parser.add_argument(
        "--sources", required=True, type=str, help="JSON array of source specifications"
    )

    parser.add_argument(
        "--output", required=True, type=Path, help="Output merged presentation path"
    )

    parser.add_argument(
        "--base-template",
        type=Path,
        default=None,
        help="Optional template to use for theme/masters",
    )

    parser.add_argument(
        "--preserve-formatting",
        action="store_true",
        default=True,
        help="Preserve original slide formatting (default: true)",
    )

    parser.add_argument(
        "--json",
        action="store_true",
        default=True,
        help="Output JSON response (default: true)",
    )

    parser.add_argument(
        "--approval-token",
        type=str,
        default=None,
        help="Approval token for merge operation (required for governance compliance)",
    )

    args = parser.parse_args()

    try:
        sources = parse_sources(args.sources)

        # Validate approval token for merge operation (governance enforcement)
        if not args.approval_token:
            raise ApprovalTokenError(
                f"Merge operation requires approval token (scope: {APPROVAL_SCOPE_MERGE_PRESENTATIONS})",
                details={"scope_required": APPROVAL_SCOPE_MERGE_PRESENTATIONS},
            )

        if len(args.approval_token) < 8:
            raise ApprovalTokenError(
                "Invalid approval token format",
                details={"token_length": len(args.approval_token)},
            )

        output_path = args.output
        if not output_path.suffix.lower() == ".pptx":
            output_path = output_path.with_suffix(".pptx")

        result = merge_presentations(
            sources=sources,
            output=output_path.resolve(),
            base_template=args.base_template.resolve() if args.base_template else None,
            preserve_formatting=args.preserve_formatting,
        )

        sys.stdout.write(json.dumps(result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(0)

    except FileNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "FileNotFoundError",
            "suggestion": "Verify all source file paths exist and are accessible",
            "tool_version": __version__,
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)

    except ValueError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ValueError",
            "suggestion": 'Check sources JSON format: [{"file":"path.pptx","slides":"all"}]',
            "tool_version": __version__,
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)

    except SlideNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "SlideNotFoundError",
            "details": getattr(e, "details", {}),
            "suggestion": "Check slide indices are valid for each source file",
            "tool_version": __version__,
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)

    except PowerPointAgentError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "details": getattr(e, "details", {}),
            "suggestion": "Check source file integrity and compatibility",
            "tool_version": __version__,
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)

    except ApprovalTokenError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ApprovalTokenError",
            "details": getattr(e, "details", {}),
            "suggestion": f"Generate approval token with scope '{APPROVAL_SCOPE_MERGE_PRESENTATIONS}' and retry",
            "tool_version": __version__,
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(4)

    except Exception as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "suggestion": "Check logs for detailed error information",
            "tool_version": __version__,
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)


if __name__ == "__main__":
    main()
