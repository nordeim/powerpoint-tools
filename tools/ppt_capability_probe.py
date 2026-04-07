#!/usr/bin/env python3
"""
PowerPoint Capability Probe Tool v3.1.0
Detect and report presentation template capabilities, layouts, and theme properties.

This tool provides comprehensive introspection of PowerPoint presentations to detect:
- Available layouts and their placeholders (with accurate runtime positions)
- Slide dimensions and aspect ratios
- Theme colors and fonts (using proper font scheme API)
- Template capabilities (footer support, slide numbers, dates)
- Multiple master slide support

Critical for AI agents and automation workflows to understand template capabilities
before generating content.

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.0

Usage:
    # Basic probe (essential info)
    uv run tools/ppt_capability_probe.py --file template.pptx --json
    
    # Deep probe (accurate positions via transient instantiation)
    uv run tools/ppt_capability_probe.py --file template.pptx --deep --json
    
    # Human-friendly summary
    uv run tools/ppt_capability_probe.py --file template.pptx --summary

Exit Codes:
    0: Success
    1: Error occurred

Design Principles:
    - Read-only operation (atomic, no file mutation)
    - JSON-first output with consistent contract
    - Accurate data via transient slide instantiation
    - Graceful degradation for missing features
    - Performance-optimized with timeout protection
"""

import sys
import os

# --- HYGIENE BLOCK START ---
# CRITICAL: Redirect stderr to /dev/null immediately to prevent library noise.
# This guarantees that `jq` or other parsers only see valid JSON on stdout.
sys.stderr = open(os.devnull, 'w')
# --- HYGIENE BLOCK END ---

import json
import argparse
import hashlib
import uuid
import time
import logging
import importlib.metadata
from pathlib import Path
from typing import Dict, Any, List, Optional, Tuple
from datetime import datetime

# Configure logging to null handler
logging.basicConfig(level=logging.CRITICAL)

# Add parent directory to path for core import
sys.path.insert(0, str(Path(__file__).parent.parent))

# ============================================================================
# CONSTANTS
# ============================================================================

__version__ = "3.1.0"
SCHEMA_VERSION = "capability_probe.v3.1.0"

# ============================================================================
# IMPORTS WITH ERROR HANDLING
# ============================================================================

try:
    from pptx import Presentation
    from pptx.enum.shapes import PP_PLACEHOLDER
except ImportError as e:
    error_output = {
        "status": "error",
        "error": f"python-pptx not installed: {e}",
        "error_type": "ImportError",
        "suggestion": "Install python-pptx: pip install python-pptx"
    }
    sys.stdout.write(json.dumps(error_output, indent=2) + "\n")
    sys.exit(1)

try:
    from core.powerpoint_agent_core import PowerPointAgentError
except ImportError:
    class PowerPointAgentError(Exception):
        """Fallback exception class if core not available."""
        pass

try:
    from core.strict_validator import validate_against_schema
    STRICT_VALIDATION_AVAILABLE = True
except ImportError:
    STRICT_VALIDATION_AVAILABLE = False
    def validate_against_schema(data, schema_path):
        pass


# ============================================================================
# LIBRARY VERSION DETECTION
# ============================================================================

def get_library_versions() -> Dict[str, str]:
    """
    Detect versions of key libraries.
    
    Returns:
        Dict mapping library name to version string
    """
    versions = {}
    
    try:
        versions["python-pptx"] = importlib.metadata.version("python-pptx")
    except importlib.metadata.PackageNotFoundError:
        versions["python-pptx"] = "not_installed"
    except Exception:
        versions["python-pptx"] = "unknown"
    
    try:
        versions["Pillow"] = importlib.metadata.version("Pillow")
    except importlib.metadata.PackageNotFoundError:
        versions["Pillow"] = "not_installed"
    except Exception:
        versions["Pillow"] = "unknown"
    
    return versions


# ============================================================================
# PLACEHOLDER TYPE MAPPING
# ============================================================================

def build_placeholder_type_map() -> Dict[int, str]:
    """
    Build mapping from PP_PLACEHOLDER enum values to human-readable names.
    
    Uses actual python-pptx enum values, not guessed numbers.
    
    Returns:
        Dict mapping type code to name
    """
    type_map = {}
    
    for name in dir(PP_PLACEHOLDER):
        if name.isupper() and not name.startswith('_'):
            try:
                member = getattr(PP_PLACEHOLDER, name)
                if isinstance(member, int):
                    code = member
                elif hasattr(member, 'value'):
                    code = member.value
                else:
                    continue
                if code is not None:
                    type_map[int(code)] = name
            except (AttributeError, TypeError, ValueError):
                pass
    
    return type_map


PLACEHOLDER_TYPE_MAP = build_placeholder_type_map()


def get_placeholder_type_name(ph_type_code: int) -> str:
    """
    Get human-readable name for placeholder type code.
    
    Args:
        ph_type_code: Numeric type code from placeholder
        
    Returns:
        Type name or UNKNOWN_X if not recognized
    """
    return PLACEHOLDER_TYPE_MAP.get(ph_type_code, f"UNKNOWN_{ph_type_code}")


# ============================================================================
# FILE UTILITIES
# ============================================================================

def calculate_file_checksum(filepath: Path) -> str:
    """
    Calculate MD5 checksum of file to verify no mutation.
    
    Args:
        filepath: Path to file
        
    Returns:
        Hex digest of file contents
    """
    md5 = hashlib.md5()
    with open(filepath, 'rb') as f:
        for chunk in iter(lambda: f.read(8192), b''):
            md5.update(chunk)
    return md5.hexdigest()


# ============================================================================
# COLOR UTILITIES
# ============================================================================

def rgb_to_hex(rgb_color) -> str:
    """
    Convert RGBColor to hex string.
    
    Args:
        rgb_color: RGBColor object from python-pptx
        
    Returns:
        Hex color string like "#0070C0"
    """
    try:
        return f"#{rgb_color.r:02X}{rgb_color.g:02X}{rgb_color.b:02X}"
    except (AttributeError, TypeError):
        return "#000000"


# ============================================================================
# DIMENSION DETECTION
# ============================================================================

def detect_slide_dimensions(prs) -> Dict[str, Any]:
    """
    Detect slide dimensions and calculate aspect ratio.
    
    Args:
        prs: Presentation object
        
    Returns:
        Dict with width, height, aspect ratio, DPI estimate
    """
    width_inches = prs.slide_width.inches
    height_inches = prs.slide_height.inches
    
    width_emu = int(prs.slide_width)
    height_emu = int(prs.slide_height)
    
    dpi_estimate = 96
    width_pixels = int(width_inches * dpi_estimate)
    height_pixels = int(height_inches * dpi_estimate)
    
    ratio = width_inches / height_inches if height_inches > 0 else 1.0
    if abs(ratio - 16/9) < 0.01:
        aspect_ratio = "16:9"
    elif abs(ratio - 4/3) < 0.01:
        aspect_ratio = "4:3"
    else:
        from fractions import Fraction
        frac = Fraction(width_pixels, height_pixels).limit_denominator(20)
        aspect_ratio = f"{frac.numerator}:{frac.denominator}"
    
    return {
        "width_inches": round(width_inches, 2),
        "height_inches": round(height_inches, 2),
        "width_emu": width_emu,
        "height_emu": height_emu,
        "width_pixels": width_pixels,
        "height_pixels": height_pixels,
        "aspect_ratio": aspect_ratio,
        "aspect_ratio_float": round(ratio, 4),
        "dpi_estimate": dpi_estimate
    }


# ============================================================================
# PLACEHOLDER ANALYSIS
# ============================================================================

def analyze_placeholder(shape, slide_width: float, slide_height: float, instantiated: bool = False) -> Dict[str, Any]:
    """
    Analyze a single placeholder and return comprehensive info.
    
    Args:
        shape: Placeholder shape to analyze
        slide_width: Slide width in inches
        slide_height: Slide height in inches
        instantiated: Whether this is from an instantiated slide (accurate) or template
        
    Returns:
        Dict with type, position, size information
    """
    ph_format = shape.placeholder_format
    ph_type = ph_format.type if hasattr(ph_format.type, '__int__') else int(ph_format.type)
    ph_type_name = get_placeholder_type_name(ph_type)
    
    try:
        left_emu = shape.left if hasattr(shape, 'left') else 0
        top_emu = shape.top if hasattr(shape, 'top') else 0
        width_emu = shape.width if hasattr(shape, 'width') else 0
        height_emu = shape.height if hasattr(shape, 'height') else 0
        
        left_inches = left_emu / 914400
        top_inches = top_emu / 914400
        width_inches = width_emu / 914400
        height_inches = height_emu / 914400
        
        left_percent = (left_inches / slide_width * 100) if slide_width > 0 else 0
        top_percent = (top_inches / slide_height * 100) if slide_height > 0 else 0
        width_percent = (width_inches / slide_width * 100) if slide_width > 0 else 0
        height_percent = (height_inches / slide_height * 100) if slide_height > 0 else 0
        
        return {
            "type": ph_type_name,
            "type_code": ph_type,
            "idx": ph_format.idx,
            "name": shape.name,
            "position_inches": {
                "left": round(left_inches, 2),
                "top": round(top_inches, 2)
            },
            "position_percent": {
                "left": f"{left_percent:.1f}%",
                "top": f"{top_percent:.1f}%"
            },
            "position_emu": {
                "left": left_emu,
                "top": top_emu
            },
            "size_inches": {
                "width": round(width_inches, 2),
                "height": round(height_inches, 2)
            },
            "size_percent": {
                "width": f"{width_percent:.1f}%",
                "height": f"{height_percent:.1f}%"
            },
            "size_emu": {
                "width": width_emu,
                "height": height_emu
            },
            "position_source": "instantiated" if instantiated else "template"
        }
    except Exception as e:
        return {
            "type": ph_type_name,
            "type_code": ph_type,
            "idx": ph_format.idx,
            "error": str(e),
            "position_source": "error"
        }


# ============================================================================
# TRANSIENT SLIDE PATTERN
# ============================================================================

def _add_transient_slide(prs, layout):
    """
    Helper to safely add and remove a transient slide for deep analysis.
    Yields the slide object, then ensures cleanup in finally block.
    
    This is the key pattern for getting accurate placeholder positions:
    template positions are theoretical until instantiated.
    
    Args:
        prs: Presentation object
        layout: Layout to instantiate
        
    Yields:
        The instantiated slide for analysis
        
    Note:
        NEVER call save() while a transient slide exists.
        The slide is automatically cleaned up, and since no save occurs,
        the file remains unmodified (atomic read guarantee).
    """
    slide = None
    added_index = -1
    try:
        slide = prs.slides.add_slide(layout)
        added_index = len(prs.slides) - 1
        yield slide
    finally:
        if added_index != -1 and added_index < len(prs.slides):
            try:
                rId = prs.slides._sldIdLst[added_index].rId
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[added_index]
            except Exception:
                # Suppress cleanup errors to avoid masking analysis errors
                # File is not saved, so transient slide disappears anyway
                pass


# ============================================================================
# LAYOUT DETECTION
# ============================================================================

def detect_layouts_with_instantiation(
    prs, 
    slide_width: float, 
    slide_height: float, 
    deep: bool, 
    warnings: List[str], 
    timeout_start: Optional[float] = None, 
    timeout_seconds: Optional[int] = None,
    max_layouts: Optional[int] = None
) -> List[Dict[str, Any]]:
    """
    Detect all layouts, optionally instantiating them for accurate positions.
    
    In deep mode, creates transient slides in-memory to get runtime positions,
    then discards them without saving (maintains atomic read guarantee).
    
    Args:
        prs: Presentation object
        slide_width: Slide width in inches
        slide_height: Slide height in inches
        deep: If True, instantiate layouts for accurate positions
        warnings: List to append warnings to
        timeout_start: Start time for timeout check
        timeout_seconds: Max seconds allowed
        max_layouts: Maximum number of layouts to analyze
        
    Returns:
        List of layout information dicts
    """
    layouts = []
    
    # Build mapping: layout partname -> master index
    master_map = {}
    try:
        for m_idx, master in enumerate(prs.slide_masters):
            for layout in master.slide_layouts:
                try:
                    key = layout.part.partname
                except AttributeError:
                    key = id(layout)
                master_map[key] = m_idx
    except Exception:
        pass

    layouts_to_process = list(prs.slide_layouts)
    total_layouts = len(layouts_to_process)
    
    if max_layouts and len(layouts_to_process) > max_layouts:
        layouts_to_process = layouts_to_process[:max_layouts]

    for idx, layout in enumerate(layouts_to_process):
        # Timeout check at each iteration
        if timeout_start and timeout_seconds:
            elapsed = time.perf_counter() - timeout_start
            if elapsed > timeout_seconds:
                warnings.append(
                    f"Probe timeout at layout {idx} ({elapsed:.1f}s > {timeout_seconds}s) - "
                    "returning partial results"
                )
                break

        # Get original index in case we sliced
        try:
            original_idx = list(prs.slide_layouts).index(layout)
        except ValueError:
            original_idx = idx

        # Determine master index
        try:
            key = layout.part.partname
        except AttributeError:
            key = id(layout)
            
        layout_info = {
            "index": idx,
            "original_index": original_idx, 
            "name": layout.name,
            "placeholder_count": len(layout.placeholders),
            "master_index": master_map.get(key)
        }
        
        if deep:
            try:
                instantiation_success = False
                for temp_slide in _add_transient_slide(prs, layout):
                    instantiation_success = True
                    
                    # Map instantiated placeholders by idx for lookup
                    instantiated_map = {}
                    for shape in temp_slide.placeholders:
                        try:
                            instantiated_map[shape.placeholder_format.idx] = shape
                        except (AttributeError, TypeError):
                            pass
                    
                    placeholders = []
                    for layout_ph in layout.placeholders:
                        try:
                            ph_idx = layout_ph.placeholder_format.idx
                            if ph_idx in instantiated_map:
                                ph_info = analyze_placeholder(
                                    instantiated_map[ph_idx], 
                                    slide_width, 
                                    slide_height, 
                                    instantiated=True
                                )
                            else:
                                ph_info = analyze_placeholder(
                                    layout_ph, 
                                    slide_width, 
                                    slide_height, 
                                    instantiated=False
                                )
                            placeholders.append(ph_info)
                        except Exception:
                            pass
                    
                    layout_info["placeholders"] = placeholders
                    layout_info["instantiation_complete"] = len(placeholders) == len(layout.placeholders)
                    layout_info["placeholder_expected"] = len(layout.placeholders)
                    layout_info["placeholder_instantiated"] = len(placeholders)

                if not instantiation_success:
                    raise Exception("Transient slide creation failed")
                
            except Exception as e:
                warnings.append(f"Could not instantiate layout '{layout.name}': {str(e)}")
                
                placeholders = []
                for shape in layout.placeholders:
                    try:
                        ph_info = analyze_placeholder(shape, slide_width, slide_height, instantiated=False)
                        placeholders.append(ph_info)
                    except Exception:
                        pass
                
                layout_info["placeholders"] = placeholders
                layout_info["instantiation_complete"] = False
                layout_info["placeholder_expected"] = len(layout.placeholders)
                layout_info["placeholder_instantiated"] = len(placeholders)
                layout_info["_warning"] = "Using template positions (instantiation failed)"
        
        # Build placeholder type summary
        placeholder_map = {}
        placeholder_types = []
        for shape in layout.placeholders:
            try:
                ph_type = shape.placeholder_format.type
                if hasattr(ph_type, '__int__'):
                    ph_type = int(ph_type)
                ph_type_name = get_placeholder_type_name(ph_type)
                
                placeholder_map[ph_type_name] = placeholder_map.get(ph_type_name, 0) + 1
                
                if ph_type_name not in placeholder_types:
                    placeholder_types.append(ph_type_name)
            except Exception:
                pass
        
        layout_info["placeholder_types"] = placeholder_types
        layout_info["placeholder_map"] = placeholder_map
        
        layouts.append(layout_info)
    
    return layouts


# ============================================================================
# THEME EXTRACTION
# ============================================================================

def extract_theme_colors(master_or_prs, warnings: List[str]) -> Dict[str, str]:
    """
    Extract theme colors from presentation or master using proper color scheme API.
    
    Args:
        master_or_prs: Presentation or SlideMaster object
        warnings: List to append warnings to
        
    Returns:
        Dict mapping color names to hex codes or scheme references
    """
    colors = {}
    
    try:
        if hasattr(master_or_prs, 'slide_masters'):
            slide_master = master_or_prs.slide_masters[0]
        else:
            slide_master = master_or_prs

        theme = getattr(slide_master, 'theme', None)
        if not theme:
            warnings.append("Theme object unavailable")
            return {}
            
        color_scheme = getattr(theme, 'theme_color_scheme', None)
        if not color_scheme:
            warnings.append("Theme color scheme unavailable")
            return {}
        
        color_attrs = [
            'accent1', 'accent2', 'accent3', 'accent4', 'accent5', 'accent6',
            'background1', 'background2', 'text1', 'text2', 'hyperlink', 'followed_hyperlink'
        ]
        
        non_rgb_found = False
        for color_name in color_attrs:
            try:
                color = getattr(color_scheme, color_name, None)
                if color:
                    if hasattr(color, 'r'):
                        colors[color_name] = rgb_to_hex(color)
                    else:
                        colors[color_name] = f"schemeColor:{color_name}"
                        non_rgb_found = True
            except Exception:
                pass
        
        if not colors:
            warnings.append("Theme color scheme unavailable or empty")
        elif non_rgb_found:
            warnings.append("Theme colors include non-RGB scheme references")
            
    except Exception as e:
        warnings.append(f"Theme color extraction failed: {str(e)}")
    
    return colors


def _font_name(font_obj) -> Optional[str]:
    """Helper to safely get typeface from font object."""
    if font_obj is None:
        return None
    return getattr(font_obj, 'typeface', str(font_obj))


def extract_theme_fonts(master_or_prs, warnings: List[str]) -> Dict[str, str]:
    """
    Extract theme fonts from presentation or master using proper font scheme API.
    
    Args:
        master_or_prs: Presentation or SlideMaster object
        warnings: List to append warnings to
        
    Returns:
        Dict with heading and body font names
    """
    fonts = {}
    fallback_used = False
    
    try:
        if hasattr(master_or_prs, 'slide_masters'):
            slide_master = master_or_prs.slide_masters[0]
        else:
            slide_master = master_or_prs

        theme = getattr(slide_master, 'theme', None)
        
        if theme:
            font_scheme = getattr(theme, 'font_scheme', None)
            if font_scheme:
                major = getattr(font_scheme, 'major_font', None)
                minor = getattr(font_scheme, 'minor_font', None)
                
                if major:
                    latin = getattr(major, 'latin', None)
                    ea = getattr(major, 'east_asian', None)
                    cs = getattr(major, 'complex_script', None)
                    
                    heading_font = _font_name(latin) or _font_name(ea) or _font_name(cs)
                    if heading_font:
                        fonts['heading'] = heading_font
                    
                    if ea and _font_name(ea):
                        fonts['heading_east_asian'] = _font_name(ea)
                    if cs and _font_name(cs):
                        fonts['heading_complex'] = _font_name(cs)
                
                if minor:
                    latin = getattr(minor, 'latin', None)
                    ea = getattr(minor, 'east_asian', None)
                    cs = getattr(minor, 'complex_script', None)
                    
                    body_font = _font_name(latin) or _font_name(ea) or _font_name(cs)
                    if body_font:
                        fonts['body'] = body_font
                    
                    if ea and _font_name(ea):
                        fonts['body_east_asian'] = _font_name(ea)
                    if cs and _font_name(cs):
                        fonts['body_complex'] = _font_name(cs)

        if not fonts:
            for shape in slide_master.shapes:
                if hasattr(shape, 'text_frame') and shape.text_frame.paragraphs:
                    for paragraph in shape.text_frame.paragraphs:
                        if paragraph.font.name and 'heading' not in fonts:
                            fonts['heading'] = paragraph.font.name
                            break
                    if 'heading' in fonts:
                        break
        
        if not fonts:
            fallback_used = True
            fonts = {"heading": "Calibri", "body": "Calibri"}
            
    except Exception as e:
        fallback_used = True
        fonts = {"heading": "Calibri", "body": "Calibri"}
        warnings.append(f"Theme font extraction failed: {str(e)}")
    
    if fallback_used and hasattr(master_or_prs, 'slide_masters'):
        warnings.append("Theme fonts unavailable - using Calibri defaults")
    
    return fonts


# ============================================================================
# CAPABILITY ANALYSIS
# ============================================================================

def analyze_capabilities(layouts: List[Dict[str, Any]], prs) -> Dict[str, Any]:
    """
    Analyze template capabilities based on detected layouts.
    
    Args:
        layouts: List of layout information dicts
        prs: Presentation object
        
    Returns:
        Dict with capability flags, layout mappings, and recommendations
    """
    has_footer = False
    has_slide_number = False
    has_date = False
    layouts_with_footer = []
    layouts_with_slide_number = []
    layouts_with_date = []
    
    footer_type_code = None
    slide_number_type_code = None
    date_type_code = None
    
    for type_code, type_name in PLACEHOLDER_TYPE_MAP.items():
        if type_name == 'FOOTER':
            footer_type_code = type_code
        elif type_name == 'SLIDE_NUMBER':
            slide_number_type_code = type_code
        elif type_name == 'DATE':
            date_type_code = type_code
    
    per_master_stats = {}
    
    for layout in layouts:
        layout_ref = {
            "index": layout['index'],
            "original_index": layout.get('original_index', layout['index']),
            "name": layout['name'],
            "master_index": layout.get('master_index')
        }
        m_idx = layout.get('master_index')
        
        if m_idx is not None:
            if m_idx not in per_master_stats:
                per_master_stats[m_idx] = {
                    "master_index": m_idx,
                    "layout_count": 0,
                    "has_footer_layouts": 0,
                    "has_slide_number_layouts": 0,
                    "has_date_layouts": 0
                }
            per_master_stats[m_idx]["layout_count"] += 1
        
        layout_has_footer = False
        layout_has_slide_number = False
        layout_has_date = False

        if 'placeholders' in layout:
            for ph in layout['placeholders']:
                if footer_type_code and ph.get('type_code') == footer_type_code:
                    has_footer = True
                    layout_has_footer = True
                    if layout_ref not in layouts_with_footer:
                        layouts_with_footer.append(layout_ref)
                
                if slide_number_type_code and ph.get('type_code') == slide_number_type_code:
                    has_slide_number = True
                    layout_has_slide_number = True
                    if layout_ref not in layouts_with_slide_number:
                        layouts_with_slide_number.append(layout_ref)
                
                if date_type_code and ph.get('type_code') == date_type_code:
                    has_date = True
                    layout_has_date = True
                    if layout_ref not in layouts_with_date:
                        layouts_with_date.append(layout_ref)
                        
        elif 'placeholder_types' in layout:
            if 'FOOTER' in layout['placeholder_types']:
                has_footer = True
                layout_has_footer = True
                layouts_with_footer.append(layout_ref)
            
            if 'SLIDE_NUMBER' in layout['placeholder_types']:
                has_slide_number = True
                layout_has_slide_number = True
                layouts_with_slide_number.append(layout_ref)
            
            if 'DATE' in layout['placeholder_types']:
                has_date = True
                layout_has_date = True
                layouts_with_date.append(layout_ref)
        
        if m_idx is not None:
            if layout_has_footer:
                per_master_stats[m_idx]["has_footer_layouts"] += 1
            if layout_has_slide_number:
                per_master_stats[m_idx]["has_slide_number_layouts"] += 1
            if layout_has_date:
                per_master_stats[m_idx]["has_date_layouts"] += 1
    
    recommendations = []
    
    if not has_footer:
        recommendations.append(
            "No footer placeholders found - ppt_set_footer.py will use text box fallback strategy"
        )
    else:
        layout_names = [l['name'] for l in layouts_with_footer]
        recommendations.append(
            f"Footer placeholders available on {len(layouts_with_footer)} layout(s): {', '.join(layout_names)}"
        )
    
    if not has_slide_number:
        recommendations.append(
            "No slide number placeholders - recommend manual text box for slide numbers"
        )
    else:
        layout_names = [l['name'] for l in layouts_with_slide_number]
        recommendations.append(
            f"Slide number placeholders available on {len(layouts_with_slide_number)} layout(s): {', '.join(layout_names)}"
        )
    
    if not has_date:
        recommendations.append(
            "No date placeholders - dates must be added manually if needed"
        )
    else:
        layout_names = [l['name'] for l in layouts_with_date]
        recommendations.append(
            f"Date placeholders available on {len(layouts_with_date)} layout(s): {', '.join(layout_names)}"
        )
    
    return {
        "has_footer_placeholders": has_footer,
        "has_slide_number_placeholders": has_slide_number,
        "has_date_placeholders": has_date,
        "layouts_with_footer": layouts_with_footer,
        "layouts_with_slide_number": layouts_with_slide_number,
        "layouts_with_date": layouts_with_date,
        "total_layouts": len(layouts),
        "total_master_slides": len(prs.slide_masters),
        "per_master": list(per_master_stats.values()),
        "footer_support_mode": "placeholder" if has_footer else "fallback_textbox",
        "slide_number_strategy": "placeholder" if has_slide_number else "textbox",
        "recommendations": recommendations
    }


# ============================================================================
# OUTPUT VALIDATION
# ============================================================================

def validate_output(result: Dict[str, Any]) -> Tuple[bool, List[str]]:
    """
    Validate probe result has all required fields.
    
    Args:
        result: Probe result dict
        
    Returns:
        Tuple of (is_valid, list of missing fields)
    """
    required_fields = [
        "status",
        "metadata",
        "metadata.file",
        "metadata.probed_at",
        "metadata.tool_version",
        "metadata.operation_id",
        "metadata.duration_ms",
        "slide_dimensions",
        "layouts",
        "theme",
        "capabilities",
        "warnings"
    ]
    
    missing = []
    
    for field_path in required_fields:
        parts = field_path.split('.')
        current = result
        
        for part in parts:
            if not isinstance(current, dict) or part not in current:
                missing.append(field_path)
                break
            current = current[part]
    
    return (len(missing) == 0, missing)


# ============================================================================
# MAIN PROBE FUNCTION
# ============================================================================

def probe_presentation(
    filepath: Path,
    deep: bool = False,
    verify_atomic: bool = True,
    max_layouts: Optional[int] = None,
    timeout_seconds: Optional[int] = None
) -> Dict[str, Any]:
    """
    Probe presentation and return comprehensive capability report.
    
    Args:
        filepath: Path to PowerPoint file
        deep: If True, perform deep analysis with transient slide instantiation
        verify_atomic: If True, verify no file mutation occurred
        max_layouts: Maximum layouts to analyze (None = all)
        timeout_seconds: Maximum seconds for analysis (None = no limit)
        
    Returns:
        Dict with complete capability report
        
    Raises:
        FileNotFoundError: If file doesn't exist
        PermissionError: If file is locked
        PowerPointAgentError: If atomic verification fails
    """
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    
    if not filepath.is_file():
        raise ValueError(f"Path is not a file: {filepath}")
    
    try:
        with open(filepath, 'rb') as f:
            f.read(1)
    except PermissionError:
        raise PermissionError(f"File is locked or permission denied: {filepath}")
    
    start_time = time.perf_counter()
    operation_id = str(uuid.uuid4())
    warnings = []
    info = []
    
    checksum_before = None
    if verify_atomic:
        checksum_before = calculate_file_checksum(filepath)
    
    prs = Presentation(str(filepath))
    
    dimensions = detect_slide_dimensions(prs)
    slide_width = dimensions['width_inches']
    slide_height = dimensions['height_inches']
    
    all_layouts = list(prs.slide_layouts)
    if max_layouts and len(all_layouts) > max_layouts:
        info.append(f"Limited analysis to first {max_layouts} of {len(all_layouts)} layouts")
    
    layouts = detect_layouts_with_instantiation(
        prs, 
        slide_width, 
        slide_height, 
        deep, 
        warnings, 
        timeout_start=start_time, 
        timeout_seconds=timeout_seconds,
        max_layouts=max_layouts
    )
    
    analysis_complete = True
    if timeout_seconds and (time.perf_counter() - start_time) > timeout_seconds:
        analysis_complete = False
    
    theme_colors = extract_theme_colors(prs, warnings)
    theme_fonts = extract_theme_fonts(prs, warnings)
    
    theme_per_master = []
    try:
        for m_idx, master in enumerate(prs.slide_masters):
            m_warnings = []
            m_colors = extract_theme_colors(master, m_warnings)
            m_fonts = extract_theme_fonts(master, m_warnings)
            theme_per_master.append({
                "master_index": m_idx,
                "colors": m_colors,
                "fonts": m_fonts
            })
    except Exception:
        pass
    
    capabilities = analyze_capabilities(layouts, prs)
    capabilities["analysis_complete"] = analysis_complete
    
    duration_ms = int((time.perf_counter() - start_time) * 1000)
    
    if verify_atomic:
        checksum_after = calculate_file_checksum(filepath)
        if checksum_before != checksum_after:
            raise PowerPointAgentError(
                "File was modified during probe operation! "
                "This should never happen (atomic read violation). "
                f"Checksum before: {checksum_before}, after: {checksum_after}"
            )
    
    masters_info = []
    try:
        for m_idx, m in enumerate(prs.slide_masters):
            masters_info.append({
                "master_index": m_idx,
                "layout_count": len(m.slide_layouts),
                "name": getattr(m, 'name', f"Master {m_idx}")
            })
    except Exception:
        pass
    
    result = {
        "status": "success",
        "metadata": {
            "file": str(filepath.resolve()),
            "probed_at": datetime.now().isoformat(),
            "tool_version": __version__,
            "schema_version": SCHEMA_VERSION,
            "operation_id": operation_id,
            "deep_analysis": deep,
            "analysis_mode": "deep" if deep else "essential",
            "atomic_verified": verify_atomic,
            "duration_ms": duration_ms,
            "timeout_seconds": timeout_seconds,
            "layout_count_total": len(all_layouts),
            "layout_count_analyzed": len(layouts),
            "warnings_count": len(warnings),
            "masters": masters_info,
            "library_versions": get_library_versions(),
            "checksum": checksum_before if verify_atomic else "verification_skipped"
        },
        "slide_dimensions": dimensions,
        "layouts": layouts,
        "theme": {
            "colors": theme_colors,
            "fonts": theme_fonts,
            "per_master": theme_per_master
        },
        "capabilities": capabilities,
        "warnings": warnings,
        "info": info
    }
    
    is_valid, missing_fields = validate_output(result)
    if not is_valid:
        warnings.append(f"Output validation found missing fields: {', '.join(missing_fields)}")
    
    if STRICT_VALIDATION_AVAILABLE:
        try:
            schema_path = Path(__file__).parent.parent / "schemas" / "capability_probe.v3.1.0.schema.json"
            if schema_path.exists():
                validate_against_schema(result, str(schema_path))
        except FileNotFoundError:
            info.append("Schema file not found - strict validation skipped")
        except Exception as e:
            warnings.append(f"Strict schema validation warning: {str(e)}")
    
    return result


# ============================================================================
# HUMAN-READABLE SUMMARY
# ============================================================================

def format_summary(probe_result: Dict[str, Any]) -> str:
    """
    Format probe result as human-readable summary.
    
    Args:
        probe_result: Result from probe_presentation()
        
    Returns:
        Formatted string summary
    """
    lines = []
    
    lines.append("═══════════════════════════════════════════════════════════════")
    lines.append(f"PowerPoint Capability Probe Report v{__version__}")
    lines.append("═══════════════════════════════════════════════════════════════")
    lines.append("")
    
    meta = probe_result['metadata']
    lines.append(f"File: {meta['file']}")
    lines.append(f"Probed: {meta['probed_at']}")
    lines.append(f"Operation ID: {meta['operation_id']}")
    lines.append(f"Analysis Mode: {'Deep (instantiated positions)' if meta['deep_analysis'] else 'Essential (template positions)'}")
    lines.append(f"Duration: {meta['duration_ms']}ms")
    lines.append(f"Atomic Verified: {'✓' if meta['atomic_verified'] else '✗'}")
    lines.append("")
    
    if meta.get('library_versions'):
        lines.append("Library Versions:")
        for lib, ver in meta['library_versions'].items():
            lines.append(f"  {lib}: {ver}")
        lines.append("")
    
    dims = probe_result['slide_dimensions']
    lines.append("Slide Dimensions:")
    lines.append(f"  Size: {dims['width_inches']}\" × {dims['height_inches']}\" ({dims['width_pixels']}×{dims['height_pixels']}px)")
    lines.append(f"  Aspect Ratio: {dims['aspect_ratio']}")
    lines.append(f"  DPI Estimate: {dims['dpi_estimate']}")
    lines.append("")
    
    caps = probe_result['capabilities']
    lines.append("Template Capabilities:")
    lines.append(f"  ✓ Total Layouts: {caps['total_layouts']}")
    lines.append(f"  ✓ Master Slides: {caps['total_master_slides']}")
    lines.append(f"  {'✓' if caps['has_footer_placeholders'] else '✗'} Footer Placeholders: {len(caps['layouts_with_footer'])} layout(s)")
    lines.append(f"  {'✓' if caps['has_slide_number_placeholders'] else '✗'} Slide Number Placeholders: {len(caps['layouts_with_slide_number'])} layout(s)")
    lines.append(f"  {'✓' if caps['has_date_placeholders'] else '✗'} Date Placeholders: {len(caps['layouts_with_date'])} layout(s)")
    lines.append("")

    if 'per_master' in caps and caps['per_master']:
        lines.append("Master Slides:")
        for m in caps['per_master']:
            lines.append(f"  Master {m['master_index']}: {m['layout_count']} layouts")
            footer = 'Yes' if m['has_footer_layouts'] else 'No'
            slide_num = 'Yes' if m['has_slide_number_layouts'] else 'No'
            date = 'Yes' if m['has_date_layouts'] else 'No'
            lines.append(f"    Footer: {footer} | Slide #: {slide_num} | Date: {date}")
        lines.append("")
    
    lines.append("Available Layouts:")
    for layout in probe_result['layouts']:
        ph_count = layout['placeholder_count']
        display_idx = layout.get('original_index', layout['index'])
        lines.append(f"  [{display_idx}] {layout['name']} ({ph_count} placeholder{'s' if ph_count != 1 else ''})")
        
        if 'placeholder_types' in layout and layout['placeholder_types']:
            types_str = ', '.join(layout['placeholder_types'])
            lines.append(f"      Types: {types_str}")
    lines.append("")
    
    theme = probe_result['theme']
    if theme.get('fonts'):
        lines.append("Theme Fonts:")
        for key, value in theme['fonts'].items():
            if not key.startswith('_'):
                lines.append(f"  {key.replace('_', ' ').title()}: {value}")
        lines.append("")
    
    if theme.get('colors'):
        color_count = len([k for k in theme['colors'].keys() if not k.startswith('_')])
        lines.append(f"Theme Colors: {color_count} defined")
        lines.append("")
    
    if caps.get('recommendations'):
        lines.append("Recommendations:")
        for rec in caps['recommendations']:
            lines.append(f"  • {rec}")
        lines.append("")
    
    if probe_result.get('warnings'):
        lines.append("⚠️  Warnings:")
        for warning in probe_result['warnings']:
            lines.append(f"  • {warning}")
        lines.append("")
    
    if probe_result.get('info'):
        lines.append("ℹ️  Information:")
        for info_msg in probe_result['info']:
            lines.append(f"  • {info_msg}")
        lines.append("")
    
    lines.append("═══════════════════════════════════════════════════════════════")
    
    return "\n".join(lines)


# ============================================================================
# CLI ENTRY POINT
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description=f"Probe PowerPoint presentation capabilities (v{__version__})",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Basic probe (essential info, fast)
  uv run tools/ppt_capability_probe.py --file template.pptx --json
  
  # Deep probe (accurate positions via transient instantiation)
  uv run tools/ppt_capability_probe.py --file template.pptx --deep --json
  
  # Human-friendly summary
  uv run tools/ppt_capability_probe.py --file template.pptx --summary
  
  # Skip atomic verification for speed
  uv run tools/ppt_capability_probe.py --file template.pptx --no-verify-atomic --json
  
  # Large template with layout limit
  uv run tools/ppt_capability_probe.py --file big_template.pptx --max-layouts 20 --json

Version: """ + __version__
    )
    
    parser.add_argument(
        '--file',
        required=True,
        type=Path,
        help='PowerPoint file to probe'
    )
    
    parser.add_argument(
        '--deep',
        action='store_true',
        help='Perform deep analysis with transient slide instantiation for accurate positions (slower)'
    )
    
    parser.add_argument(
        '--summary',
        action='store_true',
        help='Output human-friendly summary instead of JSON'
    )
    
    parser.add_argument(
        '--verify-atomic',
        action='store_true',
        default=True,
        dest='verify_atomic',
        help='Verify no file mutation occurred (default: true)'
    )
    
    parser.add_argument(
        '--no-verify-atomic',
        action='store_false',
        dest='verify_atomic',
        help='Skip atomic verification (faster, less safe)'
    )
    
    parser.add_argument(
        '--max-layouts',
        type=int,
        help='Maximum layouts to analyze (for large templates)'
    )

    parser.add_argument(
        '--timeout',
        type=int,
        default=30,
        help='Timeout in seconds for analysis (default: 30)'
    )
    
    parser.add_argument(
        '--json',
        action='store_true',
        dest='output_json',
        help='Output JSON format (default if --summary not used)'
    )
    
    args = parser.parse_args()
    
    if not args.summary and not args.output_json:
        args.output_json = True
        
    if args.summary and args.output_json:
        error_output = {
            "status": "error",
            "error": "Cannot use both --summary and --json",
            "error_type": "ArgumentError"
        }
        sys.stdout.write(json.dumps(error_output, indent=2) + "\n")
        sys.exit(1)
    
    try:
        result = probe_presentation(
            filepath=args.file,
            deep=args.deep,
            verify_atomic=args.verify_atomic,
            max_layouts=args.max_layouts,
            timeout_seconds=args.timeout
        )
        
        if args.summary:
            print(format_summary(result))
        else:
            sys.stdout.write(json.dumps(result, indent=2) + "\n")
        
        sys.exit(0)
        
    except Exception as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "suggestion": "Check file path and permissions",
            "metadata": {
                "file": str(args.file) if args.file else None,
                "tool_version": __version__,
                "operation_id": str(uuid.uuid4()),
                "probed_at": datetime.now().isoformat()
            }
        }
        
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
