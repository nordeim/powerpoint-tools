#!/usr/bin/env python3
"""
PowerPoint Update Chart Data Tool v3.1.0
Update the data of an existing chart

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.0

Usage:
    uv run tools/ppt_update_chart_data.py --file deck.pptx --slide 0 --chart 0 --data new_data.json --json

Exit Codes:
    0: Success
    1: Error occurred (check error_type in JSON for details)

⚠️ LIMITATION WARNING:
    python-pptx has LIMITED chart update support. The replace_data() method
    may fail if the new data schema doesn't match the original chart exactly.
    
    If update fails, consider the alternative approach:
    1. Delete the existing chart: ppt_remove_shape.py
    2. Add a new chart with new data: ppt_add_chart.py
"""

import sys
import os

# --- HYGIENE BLOCK START ---
# CRITICAL: Redirect stderr to /dev/null to prevent library noise from corrupting JSON output
sys.stderr = open(os.devnull, 'w')
# --- HYGIENE BLOCK END ---

import json
import argparse
from pathlib import Path
from typing import Dict, Any, Optional

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import (
    PowerPointAgent, 
    PowerPointAgentError, 
    SlideNotFoundError
)

__version__ = "3.1.0"

# Import CategoryChartData safely
try:
    from pptx.chart.data import CategoryChartData
    CHART_DATA_AVAILABLE = True
except ImportError:
    CHART_DATA_AVAILABLE = False
    CategoryChartData = None


def update_chart_data(
    filepath: Path,
    slide_index: int,
    chart_index: int,
    data: Dict[str, Any]
) -> Dict[str, Any]:
    """
    Update the data of an existing chart.
    
    Replaces the chart's data with new categories and series values.
    The new data must be compatible with the existing chart type.
    
    ⚠️ LIMITATION: python-pptx's replace_data() may fail if the new
    data structure doesn't match the original. If this fails, consider
    deleting the chart and creating a new one.
    
    Args:
        filepath: Path to the PowerPoint file to modify
        slide_index: Index of the slide containing the chart (0-based)
        chart_index: Index of the chart on the slide (0-based)
        data: New chart data dict with 'categories' and 'series' keys
        
    Returns:
        Dict containing:
            - status: "success"
            - file: Absolute path to modified file
            - slide_index: Index of the slide
            - chart_index: Index of the chart
            - categories: Number of categories
            - series: Number of data series
            - data_points: Total data points updated
            - presentation_version_before: State hash before update
            - presentation_version_after: State hash after update
            - tool_version: Version of this tool
            
    Raises:
        FileNotFoundError: If file doesn't exist
        SlideNotFoundError: If slide index is out of range
        ValueError: If data format is invalid or chart not found
        RuntimeError: If chart data update fails (python-pptx limitation)
        
    Example:
        >>> data = {
        ...     "categories": ["Q1", "Q2", "Q3"],
        ...     "series": [{"name": "Sales", "values": [100, 150, 200]}]
        ... }
        >>> result = update_chart_data(
        ...     filepath=Path("presentation.pptx"),
        ...     slide_index=1,
        ...     chart_index=0,
        ...     data=data
        ... )
        >>> print(result["data_points"])
        3
    """
    # Validate file exists
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    
    # Validate data structure
    if "categories" not in data:
        raise ValueError(
            "Data must contain 'categories' key. "
            "Example: {\"categories\": [\"A\", \"B\"], \"series\": [...]}"
        )
    
    if "series" not in data or not data["series"]:
        raise ValueError(
            "Data must contain at least one series. "
            "Example: {\"series\": [{\"name\": \"Sales\", \"values\": [10, 20]}]}"
        )
    
    # Validate series data
    cat_len = len(data["categories"])
    for i, series in enumerate(data["series"]):
        if "name" not in series:
            raise ValueError(f"Series {i} missing 'name' key")
        if "values" not in series:
            raise ValueError(f"Series {i} missing 'values' key")
        if len(series["values"]) != cat_len:
            raise ValueError(
                f"Series '{series['name']}' has {len(series['values'])} values, "
                f"but there are {cat_len} categories. Counts must match."
            )
    
    # Check if CategoryChartData is available
    if not CHART_DATA_AVAILABLE:
        raise RuntimeError(
            "pptx.chart.data.CategoryChartData not available. "
            "Ensure python-pptx is properly installed."
        )
    
    with PowerPointAgent(filepath) as agent:
        agent.open(filepath)
        
        # Capture version BEFORE update
        info_before = agent.get_presentation_info()
        version_before = info_before.get("presentation_version")
        
        # Validate slide index
        total_slides = agent.get_slide_count()
        if not 0 <= slide_index < total_slides:
            raise SlideNotFoundError(
                f"Slide index {slide_index} out of range (0-{total_slides - 1})",
                details={
                    "requested_index": slide_index,
                    "available_slides": total_slides
                }
            )
        
        # NOTE: Direct prs access required for chart data manipulation
        # python-pptx requires direct access to chart objects for replace_data()
        slide = agent.prs.slides[slide_index]
        
        # Find charts on slide
        charts = [shape for shape in slide.shapes if shape.has_chart]
        
        if not charts:
            raise ValueError(
                f"No charts found on slide {slide_index}. "
                "Use ppt_add_chart.py to create a chart first."
            )
        
        if not 0 <= chart_index < len(charts):
            raise ValueError(
                f"Chart index {chart_index} out of range. "
                f"Slide has {len(charts)} chart(s) (indices 0-{len(charts) - 1})."
            )
        
        chart_shape = charts[chart_index]
        chart = chart_shape.chart
        
        # Create new chart data
        chart_data = CategoryChartData()
        chart_data.categories = data["categories"]
        
        for series in data["series"]:
            chart_data.add_series(series["name"], series["values"])
        
        # Attempt to replace data
        try:
            chart.replace_data(chart_data)
        except Exception as e:
            raise RuntimeError(
                f"Failed to update chart data: {e}. "
                "This may be due to python-pptx limitations with complex charts. "
                "Consider deleting the chart (ppt_remove_shape.py) and "
                "creating a new one (ppt_add_chart.py) instead."
            )
        
        # Save changes
        agent.save()
        
        # Capture version AFTER update
        info_after = agent.get_presentation_info()
        version_after = info_after.get("presentation_version")
    
    return {
        "status": "success",
        "file": str(filepath.resolve()),
        "slide_index": slide_index,
        "chart_index": chart_index,
        "categories": len(data["categories"]),
        "series": len(data["series"]),
        "data_points": sum(len(s["values"]) for s in data["series"]),
        "presentation_version_before": version_before,
        "presentation_version_after": version_after,
        "tool_version": __version__
    }


def main():
    parser = argparse.ArgumentParser(
        description="Update PowerPoint chart data",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
⚠️ LIMITATION WARNING:
  python-pptx has LIMITED chart update support. The replace_data()
  method may fail if the new data doesn't match the original chart.
  
  If this tool fails, use the alternative approach:
  1. Get chart position: ppt_get_slide_info.py
  2. Delete chart: ppt_remove_shape.py (with approval token)
  3. Create new chart: ppt_add_chart.py

Examples:
  # Update from JSON file
  uv run tools/ppt_update_chart_data.py \\
    --file presentation.pptx \\
    --slide 1 \\
    --chart 0 \\
    --data updated_data.json \\
    --json
  
  # Update with inline data
  uv run tools/ppt_update_chart_data.py \\
    --file presentation.pptx \\
    --slide 0 \\
    --chart 0 \\
    --data-string '{"categories":["Q1","Q2","Q3"],"series":[{"name":"Sales","values":[100,150,200]}]}' \\
    --json

Data Format (JSON):
{
  "categories": ["Q1", "Q2", "Q3", "Q4"],
  "series": [
    {"name": "Revenue", "values": [100, 120, 140, 160]},
    {"name": "Costs", "values": [80, 90, 100, 110]}
  ]
}

Requirements:
  - Number of values in each series must match number of categories
  - Each series must have 'name' and 'values' keys
  - Data structure should match original chart type

Finding Charts:
  Use ppt_get_slide_info.py to identify charts:
  uv run tools/ppt_get_slide_info.py --file deck.pptx --slide 0 --json

Common Issues:
  - "Failed to update chart data": Schema mismatch
    Solution: Delete and recreate the chart
  
  - "No charts found": Slide has no charts
    Solution: Use ppt_add_chart.py to create one
  
  - Series count mismatch may cause issues
    Solution: Match the original number of series

Output Format:
  {
    "status": "success",
    "file": "/path/to/presentation.pptx",
    "slide_index": 1,
    "chart_index": 0,
    "categories": 4,
    "series": 2,
    "data_points": 8,
    "presentation_version_before": "a1b2c3d4...",
    "presentation_version_after": "e5f6g7h8...",
    "tool_version": "3.1.0"
  }
        """
    )
    
    parser.add_argument(
        '--file', 
        required=True, 
        type=Path, 
        help='PowerPoint file path'
    )
    parser.add_argument(
        '--slide', 
        required=True, 
        type=int, 
        help='Slide index (0-based)'
    )
    parser.add_argument(
        '--chart', 
        required=True, 
        type=int, 
        help='Chart index on slide (0-based)'
    )
    parser.add_argument(
        '--data', 
        type=Path, 
        help='JSON file with chart data'
    )
    parser.add_argument(
        '--data-string',
        help='Inline JSON data string'
    )
    parser.add_argument(
        '--json', 
        action='store_true', 
        default=True, 
        help='Output JSON response (default: true)'
    )
    
    args = parser.parse_args()
    
    try:
        # Load chart data
        if args.data:
            if not args.data.exists():
                raise FileNotFoundError(f"Data file not found: {args.data}")
            with open(args.data, 'r') as f:
                data_content = json.load(f)
        elif args.data_string:
            try:
                data_content = json.loads(args.data_string)
            except json.JSONDecodeError as e:
                raise ValueError(f"Invalid JSON in --data-string: {e}")
        else:
            raise ValueError("Either --data or --data-string is required")
        
        result = update_chart_data(
            filepath=args.file, 
            slide_index=args.slide, 
            chart_index=args.chart,
            data=data_content
        )
        sys.stdout.write(json.dumps(result, indent=2) + "\n")
        sys.exit(0)
        
    except FileNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "FileNotFoundError",
            "suggestion": "Verify file paths exist and are accessible"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except SlideNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "SlideNotFoundError",
            "details": getattr(e, 'details', {}),
            "suggestion": "Use ppt_get_info.py to check available slide indices"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except ValueError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ValueError",
            "suggestion": "Check data format and chart index"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except RuntimeError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "RuntimeError",
            "suggestion": "Consider deleting chart and creating new one with ppt_add_chart.py"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except PowerPointAgentError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "details": getattr(e, 'details', {})
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except Exception as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
