#!/usr/bin/env python3
"""
PowerPoint Crop Image Tool v3.1.0
Crop an existing image on a slide by trimming edges

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.0

Usage:
    uv run tools/ppt_crop_image.py --file deck.pptx --slide 0 --shape 1 --left 0.1 --right 0.1 --json

Exit Codes:
    0: Success
    1: Error occurred (check error_type in JSON for details)

Notes:
    Crop values are percentages of the original image size (0.0 to 1.0).
    For example, --left 0.1 trims 10% from the left edge.
"""

import sys
import os

# --- HYGIENE BLOCK START ---
# CRITICAL: Redirect stderr to /dev/null to prevent library noise from corrupting JSON output
sys.stderr = open(os.devnull, 'w')
# --- HYGIENE BLOCK END ---

import json
import argparse
from pathlib import Path
from typing import Dict, Any

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import (
    PowerPointAgent, 
    PowerPointAgentError, 
    SlideNotFoundError
)

# Import MSO_SHAPE_TYPE safely
try:
    from pptx.enum.shapes import MSO_SHAPE_TYPE
except ImportError:
    # Fallback if pptx not directly importable
    MSO_SHAPE_TYPE = None

__version__ = "3.1.0"


# Define ShapeNotFoundError if not available in core
try:
    from core.powerpoint_agent_core import ShapeNotFoundError
except ImportError:
    class ShapeNotFoundError(PowerPointAgentError):
        """Exception raised when shape is not found."""
        def __init__(self, message: str, details: Dict = None):
            self.message = message
            self.details = details or {}
            super().__init__(message)


def crop_image(
    filepath: Path,
    slide_index: int,
    shape_index: int,
    left: float = 0.0,
    right: float = 0.0,
    top: float = 0.0,
    bottom: float = 0.0
) -> Dict[str, Any]:
    """
    Crop an image on a slide by trimming edges.
    
    Applies crop values to an existing image shape. Crop values represent
    the percentage of the original image to trim from each edge.
    
    Args:
        filepath: Path to the PowerPoint file to modify
        slide_index: Index of the slide containing the image (0-based)
        shape_index: Index of the image shape to crop (0-based)
        left: Percentage to crop from left edge (0.0-1.0, default: 0.0)
        right: Percentage to crop from right edge (0.0-1.0, default: 0.0)
        top: Percentage to crop from top edge (0.0-1.0, default: 0.0)
        bottom: Percentage to crop from bottom edge (0.0-1.0, default: 0.0)
        
    Returns:
        Dict containing:
            - status: "success"
            - file: Absolute path to modified file
            - slide_index: Index of the slide
            - shape_index: Index of the cropped shape
            - crop_applied: Dict with applied crop values
            - presentation_version_before: State hash before crop
            - presentation_version_after: State hash after crop
            - tool_version: Version of this tool
            
    Raises:
        FileNotFoundError: If the PowerPoint file doesn't exist
        SlideNotFoundError: If the slide index is out of range
        ShapeNotFoundError: If the shape index is out of range
        ValueError: If crop values are invalid or shape is not an image
        
    Example:
        >>> result = crop_image(
        ...     filepath=Path("presentation.pptx"),
        ...     slide_index=0,
        ...     shape_index=1,
        ...     left=0.1,
        ...     right=0.1
        ... )
        >>> print(result["crop_applied"])
        {'left': 0.1, 'right': 0.1, 'top': 0.0, 'bottom': 0.0}
    """
    # Validate file exists
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    
    # Validate crop values
    crop_values = [left, right, top, bottom]
    for name, value in [("left", left), ("right", right), ("top", top), ("bottom", bottom)]:
        if not (0.0 <= value <= 1.0):
            raise ValueError(
                f"Crop value '{name}' must be between 0.0 and 1.0, got: {value}"
            )
    
    # Validate total crop doesn't exceed 100%
    if left + right >= 1.0:
        raise ValueError(
            f"Combined left ({left}) and right ({right}) crop cannot exceed 1.0"
        )
    if top + bottom >= 1.0:
        raise ValueError(
            f"Combined top ({top}) and bottom ({bottom}) crop cannot exceed 1.0"
        )

    with PowerPointAgent(filepath) as agent:
        agent.open(filepath)
        
        # Capture version BEFORE crop
        info_before = agent.get_presentation_info()
        version_before = info_before.get("presentation_version")
        
        # Validate slide index
        total_slides = agent.get_slide_count()
        if not 0 <= slide_index < total_slides:
            raise SlideNotFoundError(
                f"Slide index {slide_index} out of range (0-{total_slides - 1})",
                details={
                    "requested_index": slide_index,
                    "available_slides": total_slides
                }
            )
        
        # NOTE: Direct prs access required because python-pptx crop API
        # requires accessing the shape's crop properties directly.
        # This is a necessary workaround for features not exposed via agent methods.
        slide = agent.prs.slides[slide_index]
        
        # Validate shape index
        if not 0 <= shape_index < len(slide.shapes):
            raise ShapeNotFoundError(
                f"Shape index {shape_index} out of range (0-{len(slide.shapes) - 1})",
                details={
                    "requested_index": shape_index,
                    "available_shapes": len(slide.shapes)
                }
            )
        
        shape = slide.shapes[shape_index]
        
        # Validate shape is a picture
        if MSO_SHAPE_TYPE is not None:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                raise ValueError(
                    f"Shape at index {shape_index} is not an image (type: {shape.shape_type}). "
                    "Use ppt_get_slide_info.py to identify image shapes."
                )
        else:
            # Fallback check if MSO_SHAPE_TYPE not available
            if not hasattr(shape, 'crop_left'):
                raise ValueError(
                    f"Shape at index {shape_index} does not support cropping. "
                    "Ensure it is an image shape."
                )
        
        # Apply crop values (only set if > 0 to avoid unnecessary changes)
        if left > 0:
            shape.crop_left = left
        if right > 0:
            shape.crop_right = right
        if top > 0:
            shape.crop_top = top
        if bottom > 0:
            shape.crop_bottom = bottom
        
        # Save changes
        agent.save()
        
        # Capture version AFTER crop
        info_after = agent.get_presentation_info()
        version_after = info_after.get("presentation_version")
    
    return {
        "status": "success",
        "file": str(filepath.resolve()),
        "slide_index": slide_index,
        "shape_index": shape_index,
        "crop_applied": {
            "left": left,
            "right": right,
            "top": top,
            "bottom": bottom
        },
        "presentation_version_before": version_before,
        "presentation_version_after": version_after,
        "tool_version": __version__
    }


def main():
    parser = argparse.ArgumentParser(
        description="Crop an image in a PowerPoint presentation",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Crop 10%% from left and right edges
  uv run tools/ppt_crop_image.py \\
    --file presentation.pptx \\
    --slide 0 \\
    --shape 1 \\
    --left 0.1 \\
    --right 0.1 \\
    --json
  
  # Crop to focus on center (trim all edges)
  uv run tools/ppt_crop_image.py \\
    --file deck.pptx \\
    --slide 2 \\
    --shape 3 \\
    --left 0.15 \\
    --right 0.15 \\
    --top 0.1 \\
    --bottom 0.1 \\
    --json
  
  # Crop top portion only
  uv run tools/ppt_crop_image.py \\
    --file presentation.pptx \\
    --slide 0 \\
    --shape 0 \\
    --top 0.2 \\
    --json

Crop Values:
  - Values are percentages of original image size (0.0 to 1.0)
  - 0.0 = no crop, 0.1 = 10%% crop, 0.5 = 50%% crop
  - Combined opposite edges (left+right or top+bottom) must be < 1.0

Finding Images:
  Use ppt_get_slide_info.py to identify image shape indices:
  uv run tools/ppt_get_slide_info.py --file deck.pptx --slide 0 --json

Output Format:
  {
    "status": "success",
    "file": "/path/to/presentation.pptx",
    "slide_index": 0,
    "shape_index": 1,
    "crop_applied": {
      "left": 0.1,
      "right": 0.1,
      "top": 0.0,
      "bottom": 0.0
    },
    "presentation_version_before": "a1b2c3d4...",
    "presentation_version_after": "e5f6g7h8...",
    "tool_version": "3.1.0"
  }
        """
    )
    
    parser.add_argument(
        '--file', 
        required=True, 
        type=Path, 
        help='PowerPoint file path'
    )
    parser.add_argument(
        '--slide', 
        required=True, 
        type=int, 
        help='Slide index (0-based)'
    )
    parser.add_argument(
        '--shape', 
        required=True, 
        type=int, 
        help='Shape index of image to crop (0-based)'
    )
    parser.add_argument(
        '--left', 
        type=float, 
        default=0.0, 
        help='Crop percentage from left edge (0.0-1.0, default: 0.0)'
    )
    parser.add_argument(
        '--right', 
        type=float, 
        default=0.0, 
        help='Crop percentage from right edge (0.0-1.0, default: 0.0)'
    )
    parser.add_argument(
        '--top', 
        type=float, 
        default=0.0, 
        help='Crop percentage from top edge (0.0-1.0, default: 0.0)'
    )
    parser.add_argument(
        '--bottom', 
        type=float, 
        default=0.0, 
        help='Crop percentage from bottom edge (0.0-1.0, default: 0.0)'
    )
    parser.add_argument(
        '--json', 
        action='store_true', 
        default=True, 
        help='Output JSON response (default: true)'
    )
    
    args = parser.parse_args()
    
    try:
        result = crop_image(
            filepath=args.file, 
            slide_index=args.slide, 
            shape_index=args.shape,
            left=args.left,
            right=args.right,
            top=args.top,
            bottom=args.bottom
        )
        sys.stdout.write(json.dumps(result, indent=2) + "\n")
        sys.exit(0)
        
    except FileNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "FileNotFoundError",
            "suggestion": "Verify the file path exists and is accessible"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except SlideNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "SlideNotFoundError",
            "details": getattr(e, 'details', {}),
            "suggestion": "Use ppt_get_info.py to check available slide indices"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except ShapeNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ShapeNotFoundError",
            "details": getattr(e, 'details', {}),
            "suggestion": "Use ppt_get_slide_info.py to check available shape indices"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except ValueError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ValueError",
            "suggestion": "Check crop values (0.0-1.0) and ensure shape is an image"
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except PowerPointAgentError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "details": getattr(e, 'details', {})
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)
        
    except Exception as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.exit(1)


if __name__ == "__main__":
    main()
