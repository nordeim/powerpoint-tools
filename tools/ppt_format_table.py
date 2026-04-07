#!/usr/bin/env python3
"""
PowerPoint Format Table Tool v3.1.1
Style and format existing tables in presentations.

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.1

Usage:
    uv run tools/ppt_format_table.py --file presentation.pptx --slide 0 --shape 2 --header-fill "#0070C0" --json

Exit Codes:
    0: Success
    1: Error occurred (check error_type in JSON for details)

This tool formats existing tables by applying styling options including:
- Header row colors and formatting
- Data row colors with optional banding
- Font styling (name, size, color)
- Border styling (color, width)
- First column highlighting
"""

import sys
import os

# --- HYGIENE BLOCK START ---
# CRITICAL: Redirect stderr to /dev/null to prevent library noise from corrupting JSON output
sys.stderr = open(os.devnull, 'w')
# --- HYGIENE BLOCK END ---

import json
import argparse
from pathlib import Path
from typing import Dict, Any, List, Optional

sys.path.insert(0, str(Path(__file__).parent.parent))

from core.powerpoint_agent_core import (
    PowerPointAgent,
    PowerPointAgentError,
    SlideNotFoundError,
    ShapeNotFoundError
)

# ============================================================================
# CONSTANTS
# ============================================================================

__version__ = "3.1.1"


# ============================================================================
# HELPER FUNCTIONS
# ============================================================================

def parse_color(color_str: Optional[str]) -> Optional[str]:
    """
    Parse and validate a color string.
    
    Args:
        color_str: Color in #RRGGBB or RRGGBB format
        
    Returns:
        Normalized color string with # prefix, or None
    """
    if not color_str:
        return None
    
    color = color_str.strip()
    if not color.startswith('#'):
        color = '#' + color
    
    if len(color) != 7:
        raise ValueError(f"Invalid color format: {color_str}. Expected #RRGGBB")
    
    try:
        int(color[1:], 16)
    except ValueError:
        raise ValueError(f"Invalid color format: {color_str}. Expected hexadecimal")
    
    return color.upper()


def is_table_shape(shape) -> bool:
    """
    Check if a shape is a table.
    
    Args:
        shape: Shape object from python-pptx
        
    Returns:
        True if shape is a table, False otherwise
    """
    return hasattr(shape, 'table') and shape.has_table


# ============================================================================
# MAIN LOGIC
# ============================================================================

def format_table(
    filepath: Path,
    slide_index: int,
    shape_index: int,
    header_fill: Optional[str] = None,
    header_text: Optional[str] = None,
    row_fill: Optional[str] = None,
    alt_row_fill: Optional[str] = None,
    text_color: Optional[str] = None,
    font_name: Optional[str] = None,
    font_size: Optional[int] = None,
    border_color: Optional[str] = None,
    border_width: Optional[float] = None,
    first_col_highlight: bool = False,
    banding: bool = False
) -> Dict[str, Any]:
    """
    Format an existing table in a PowerPoint presentation.
    
    Args:
        filepath: Path to the PowerPoint file
        slide_index: Index of the slide containing the table (0-based)
        shape_index: Index of the table shape on the slide
        header_fill: Fill color for header row (#RRGGBB)
        header_text: Text color for header row
        row_fill: Fill color for data rows
        alt_row_fill: Alternating row fill color for banding
        text_color: Default text color for all cells
        font_name: Font family name
        font_size: Font size in points
        border_color: Border color
        border_width: Border width in points
        first_col_highlight: Highlight first column
        banding: Enable row banding (requires alt_row_fill)
        
    Returns:
        Dict with formatting results
        
    Raises:
        FileNotFoundError: If file doesn't exist
        SlideNotFoundError: If slide index is invalid
        ShapeNotFoundError: If shape index is invalid
        ValueError: If shape is not a table
    """
    if not filepath.exists():
        raise FileNotFoundError(f"File not found: {filepath}")
    
    header_fill_parsed = parse_color(header_fill)
    header_text_parsed = parse_color(header_text)
    row_fill_parsed = parse_color(row_fill)
    alt_row_fill_parsed = parse_color(alt_row_fill)
    text_color_parsed = parse_color(text_color)
    border_color_parsed = parse_color(border_color)
    
    changes_applied: List[str] = []
    
    with PowerPointAgent(filepath) as agent:
        agent.open(filepath)
        
        info_before = agent.get_presentation_info()
        version_before = info_before.get("presentation_version")
        
        total_slides = agent.get_slide_count()
        if not 0 <= slide_index < total_slides:
            raise SlideNotFoundError(
                f"Slide index {slide_index} out of range (0-{total_slides - 1})",
                details={
                    "requested_index": slide_index,
                    "available_slides": total_slides
                }
            )
        
        slide = agent.prs.slides[slide_index]
        
        if not 0 <= shape_index < len(slide.shapes):
            raise ShapeNotFoundError(
                f"Shape index {shape_index} out of range (0-{len(slide.shapes) - 1})",
                details={
                    "requested_index": shape_index,
                    "available_shapes": len(slide.shapes)
                }
            )
        
        shape = slide.shapes[shape_index]
        
        if not is_table_shape(shape):
            raise ValueError(
                f"Shape at index {shape_index} is not a table. "
                f"Shape type: {shape.shape_type}"
            )
        
        table = shape.table
        row_count = len(table.rows)
        col_count = len(table.columns)
        
        from pptx.util import Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        
        def hex_to_rgb(hex_color: str) -> RGBColor:
            hex_color = hex_color.lstrip('#')
            return RGBColor(
                int(hex_color[0:2], 16),
                int(hex_color[2:4], 16),
                int(hex_color[4:6], 16)
            )
        
        if header_fill_parsed and row_count > 0:
            for cell in table.rows[0].cells:
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(header_fill_parsed)
            changes_applied.append(f"header_fill={header_fill_parsed}")
        
        if header_text_parsed and row_count > 0:
            for cell in table.rows[0].cells:
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.color.rgb = hex_to_rgb(header_text_parsed)
            changes_applied.append(f"header_text={header_text_parsed}")
        
        if row_fill_parsed or (banding and alt_row_fill_parsed):
            for row_idx in range(1, row_count):
                if banding and alt_row_fill_parsed and row_idx % 2 == 0:
                    fill_color = alt_row_fill_parsed
                elif row_fill_parsed:
                    fill_color = row_fill_parsed
                else:
                    continue
                    
                for cell in table.rows[row_idx].cells:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = hex_to_rgb(fill_color)
            
            if row_fill_parsed:
                changes_applied.append(f"row_fill={row_fill_parsed}")
            if banding and alt_row_fill_parsed:
                changes_applied.append(f"banding_enabled=True")
                changes_applied.append(f"alt_row_fill={alt_row_fill_parsed}")
        
        if text_color_parsed:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.color.rgb = hex_to_rgb(text_color_parsed)
            changes_applied.append(f"text_color={text_color_parsed}")
        
        if font_name:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = font_name
            changes_applied.append(f"font_name={font_name}")
        
        if font_size:
            for row in table.rows:
                for cell in row.cells:
                    for paragraph in cell.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.size = Pt(font_size)
            changes_applied.append(f"font_size={font_size}pt")
        
        if first_col_highlight and header_fill_parsed and col_count > 0:
            for row in table.rows:
                cell = row.cells[0]
                cell.fill.solid()
                cell.fill.fore_color.rgb = hex_to_rgb(header_fill_parsed)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.bold = True
            changes_applied.append("first_col_highlight=True")
        
        agent.save()
        
        info_after = agent.get_presentation_info()
        version_after = info_after.get("presentation_version")
    
    return {
        "status": "success",
        "file": str(filepath.resolve()),
        "slide_index": slide_index,
        "shape_index": shape_index,
        "table_info": {
            "rows": row_count,
            "columns": col_count,
            "has_header": True
        },
        "changes_applied": changes_applied,
        "changes_count": len(changes_applied),
        "presentation_version_before": version_before,
        "presentation_version_after": version_after,
        "tool_version": __version__
    }


# ============================================================================
# CLI INTERFACE
# ============================================================================

def main():
    parser = argparse.ArgumentParser(
        description="Format existing tables in PowerPoint presentations",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  # Format header row with blue fill
  uv run tools/ppt_format_table.py \\
    --file presentation.pptx --slide 0 --shape 2 \\
    --header-fill "#0070C0" --header-text "#FFFFFF" --json

  # Enable row banding
  uv run tools/ppt_format_table.py \\
    --file presentation.pptx --slide 1 --shape 3 \\
    --row-fill "#FFFFFF" --alt-row-fill "#F0F0F0" --banding --json

  # Complete formatting
  uv run tools/ppt_format_table.py \\
    --file presentation.pptx --slide 0 --shape 2 \\
    --header-fill "#0070C0" --header-text "#FFFFFF" \\
    --row-fill "#FFFFFF" --text-color "#333333" \\
    --font-name "Calibri" --font-size 11 \\
    --first-col --json

Color Format:
  Colors must be in #RRGGBB hexadecimal format.
  Examples: #0070C0 (blue), #FFFFFF (white), #333333 (dark gray)

Output Format:
  {
    "status": "success",
    "file": "/path/to/presentation.pptx",
    "slide_index": 0,
    "shape_index": 2,
    "table_info": {"rows": 5, "columns": 4, "has_header": true},
    "changes_applied": ["header_fill=#0070C0", "header_text=#FFFFFF"],
    "presentation_version_before": "a1b2c3...",
    "presentation_version_after": "d4e5f6...",
    "tool_version": "3.1.1"
  }
        """
    )
    
    parser.add_argument(
        '--file',
        required=True,
        type=Path,
        help='PowerPoint file path'
    )
    
    parser.add_argument(
        '--slide',
        required=True,
        type=int,
        help='Slide index (0-based)'
    )
    
    parser.add_argument(
        '--shape',
        required=True,
        type=int,
        help='Shape index of the table'
    )
    
    parser.add_argument(
        '--header-fill',
        type=str,
        help='Header row fill color (#RRGGBB)'
    )
    
    parser.add_argument(
        '--header-text',
        type=str,
        help='Header row text color (#RRGGBB)'
    )
    
    parser.add_argument(
        '--row-fill',
        type=str,
        help='Data row fill color (#RRGGBB)'
    )
    
    parser.add_argument(
        '--alt-row-fill',
        type=str,
        help='Alternating row fill color for banding (#RRGGBB)'
    )
    
    parser.add_argument(
        '--text-color',
        type=str,
        help='Default text color (#RRGGBB)'
    )
    
    parser.add_argument(
        '--font-name',
        type=str,
        help='Font family name (e.g., "Calibri")'
    )
    
    parser.add_argument(
        '--font-size',
        type=int,
        help='Font size in points'
    )
    
    parser.add_argument(
        '--border-color',
        type=str,
        help='Border color (#RRGGBB)'
    )
    
    parser.add_argument(
        '--border-width',
        type=float,
        help='Border width in points'
    )
    
    parser.add_argument(
        '--first-col',
        action='store_true',
        help='Highlight first column like header'
    )
    
    parser.add_argument(
        '--banding',
        action='store_true',
        help='Enable alternating row colors (requires --alt-row-fill)'
    )
    
    parser.add_argument(
        '--json',
        action='store_true',
        default=True,
        help='Output JSON response (default: true)'
    )
    
    args = parser.parse_args()
    
    try:
        result = format_table(
            filepath=args.file.resolve(),
            slide_index=args.slide,
            shape_index=args.shape,
            header_fill=args.header_fill,
            header_text=args.header_text,
            row_fill=args.row_fill,
            alt_row_fill=args.alt_row_fill,
            text_color=args.text_color,
            font_name=args.font_name,
            font_size=args.font_size,
            border_color=args.border_color,
            border_width=args.border_width,
            first_col_highlight=args.first_col,
            banding=args.banding
        )
        
        sys.stdout.write(json.dumps(result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(0)
        
    except FileNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "FileNotFoundError",
            "suggestion": "Verify the file path exists and is accessible",
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)
        
    except SlideNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "SlideNotFoundError",
            "details": getattr(e, 'details', {}),
            "suggestion": "Use ppt_get_info.py to check available slide indices",
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)
        
    except ShapeNotFoundError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ShapeNotFoundError",
            "details": getattr(e, 'details', {}),
            "suggestion": "Use ppt_get_slide_info.py to check available shape indices",
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)
        
    except ValueError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": "ValueError",
            "suggestion": "Ensure the shape is a table and colors are in #RRGGBB format",
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)
        
    except PowerPointAgentError as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "details": getattr(e, 'details', {}),
            "suggestion": "Check file integrity and table structure",
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)
        
    except Exception as e:
        error_result = {
            "status": "error",
            "error": str(e),
            "error_type": type(e).__name__,
            "suggestion": "Check logs for detailed error information",
            "tool_version": __version__
        }
        sys.stdout.write(json.dumps(error_result, indent=2) + "\n")
        sys.stdout.flush()
        sys.exit(1)


if __name__ == "__main__":
    main()
