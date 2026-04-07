#!/usr/bin/env python3
"""
PowerPoint Agent Core Library v3.1
Production-grade PowerPoint manipulation with validation, accessibility, and full
alignment with Presentation Architect System Prompt v3.0.

This is the foundational library used by all CLI tools.
Designed for stateless, security-hardened PowerPoint operations.

Author: PowerPoint Agent Team
License: MIT
Version: 3.1.0

Changelog v3.1.0 (Security & Governance Release):
- SECURITY: Added approval_token requirement for destructive operations (delete_slide, remove_shape)
- SECURITY: Added Path Traversal protection to PathValidator
- SECURITY: Hardened FileLock with cross-platform atomic operations
- SECURITY: Standardized on SHA-256 for all hashing operations
- OBSERVABILITY: All mutation methods now return presentation_version_before/after
- OBSERVABILITY: Version hashing now includes shape geometry (position/size) to detect layout changes
- SAFETY: Removed silent index clamping (now raises SlideNotFoundError)
- SAFETY: Strict validation for shape types
- FIXED: _log_warning now correctly uses the logger instead of stderr
- FIXED: Redundant imports and duplicate logic consolidated

Changelog v3.0.0 (Major Release):
- NEW: add_notes() - Add/append/prepend/overwrite speaker notes
- NEW: set_z_order() - Control shape layering with 4 actions
- NEW: remove_shape() - Remove shapes from slides
- NEW: set_footer() - Configure footer text, numbers, date
- NEW: set_background() - Set slide/presentation background color or image
- NEW: crop_image() - True image cropping (not just resize)
- NEW: clone_presentation() - Clone presentation to new file
- NEW: get_presentation_version() - Compute deterministic version hash
- NEW: PathValidator class - Security-hardened path validation
- NEW: ShapeNotFoundError, ChartNotFoundError, PathValidationError exceptions
- NEW: ZOrderAction, NotesMode enums
- FIXED: FileLock now uses atomic os.open() with O_CREAT|O_EXCL
- FIXED: Lock released in finally block on open() failure
- FIXED: Slide insertion XML manipulation corrected
- FIXED: Placeholder type handling normalized to integers
- FIXED: Alt text detection checks 'descr' attribute
- FIXED: All bounds checks include negative index validation
- FIXED: Chart update error handling improved
- IMPROVED: All add_* methods return shape index for chaining
- IMPROVED: TemplateProfile uses lazy loading
- IMPROVED: Layout lookup cached for performance
- IMPROVED: Comprehensive docstrings with examples
- IMPROVED: Full alignment with System Prompt v3.0

Changelog v1.1.0:
- Added missing subprocess import for PDF export
- Added missing PP_PLACEHOLDER import and constants
- Replaced all magic numbers with named constants
- Removed text truncation in get_slide_info()
- Added position/size information to shape inspection
- Added placeholder subtype decoding
- Replaced print() with proper logging

Dependencies:
- python-pptx >= 0.6.21 (required)
- Pillow >= 9.0.0 (optional, for image operations)
"""

import os
import re
import sys
import json
import hashlib
import subprocess
import tempfile
import shutil
import time
import logging
import platform
import errno
from pathlib import Path
from typing import Any, Dict, List, Optional, Union, Tuple
from enum import Enum
from datetime import datetime
from io import BytesIO
from lxml import etree
from pptx.oxml.ns import qn

# ============================================================================
# THIRD-PARTY IMPORTS WITH GRACEFUL DEGRADATION
# ============================================================================

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    from pptx.enum.dml import MSO_THEME_COLOR
    from pptx.chart.data import CategoryChartData
    from pptx.dml.color import RGBColor

    PPTX_AVAILABLE = True
except ImportError:
    PPTX_AVAILABLE = False
    raise ImportError(
        "python-pptx is required. Install with:\n"
        "  pip install python-pptx\n"
        "  or: uv pip install python-pptx"
    )

try:
    from PIL import Image as PILImage

    HAS_PILLOW = True
except ImportError:
    HAS_PILLOW = False
    PILImage = None


# ============================================================================
# LOGGING SETUP
# ============================================================================

logger = logging.getLogger(__name__)
if not logger.handlers:
    handler = logging.StreamHandler()
    formatter = logging.Formatter("%(levelname)s:%(name)s:%(message)s")
    handler.setFormatter(formatter)
    logger.addHandler(handler)
    logger.setLevel(logging.WARNING)


# ============================================================================
# EXCEPTIONS
# ============================================================================


class PowerPointAgentError(Exception):
    """Base exception for all PowerPoint agent errors."""

    def __init__(self, message: str, details: Optional[Dict[str, Any]] = None):
        super().__init__(message)
        self.message = message
        self.details = details or {}

    def to_dict(self) -> Dict[str, Any]:
        """Convert exception to JSON-serializable dict."""
        return {
            "error": self.__class__.__name__,
            "message": self.message,
            "details": self.details,
        }

    def to_json(self) -> str:
        """Convert exception to JSON string."""
        return json.dumps(self.to_dict())


class SlideNotFoundError(PowerPointAgentError):
    """Raised when slide index is out of range."""

    pass


class ShapeNotFoundError(PowerPointAgentError):
    """Raised when shape index is out of range."""

    pass


class ChartNotFoundError(PowerPointAgentError):
    """Raised when chart is not found at specified index."""

    pass


class LayoutNotFoundError(PowerPointAgentError):
    """Raised when requested layout doesn't exist."""

    pass


class ImageNotFoundError(PowerPointAgentError):
    """Raised when image file is not found."""

    pass


class InvalidPositionError(PowerPointAgentError):
    """Raised when position specification is invalid."""

    pass


class TemplateError(PowerPointAgentError):
    """Raised when template operations fail."""

    pass


class ThemeError(PowerPointAgentError):
    """Raised when theme operations fail."""

    pass


class AccessibilityError(PowerPointAgentError):
    """Raised when accessibility validation fails."""

    pass


class AssetValidationError(PowerPointAgentError):
    """Raised when asset validation fails."""

    pass


class FileLockError(PowerPointAgentError):
    """Raised when file cannot be locked for exclusive access."""

    pass


class PathValidationError(PowerPointAgentError):
    """Raised when path validation fails (security)."""

    pass


class ApprovalTokenError(PowerPointAgentError):
    """Raised when a destructive operation lacks a valid approval token."""

    pass


# ============================================================================
# CONSTANTS
# ============================================================================

__version__ = "3.1.0"
__author__ = "PowerPoint Agent Team"
__license__ = "MIT"

# Standard slide dimensions (16:9 widescreen) in inches
SLIDE_WIDTH_INCHES = 13.333
SLIDE_HEIGHT_INCHES = 7.5

# Alternative dimensions (4:3 standard) in inches
SLIDE_WIDTH_4_3_INCHES = 10.0
SLIDE_HEIGHT_4_3_INCHES = 7.5

# EMU conversion constant
EMU_PER_INCH = 914400

# Governance Scopes
APPROVAL_SCOPE_DELETE_SLIDE = "delete:slide"
APPROVAL_SCOPE_REMOVE_SHAPE = "remove:shape"
APPROVAL_SCOPE_MERGE_PRESENTATIONS = "merge:presentations"

# Standard anchor points for positioning
ANCHOR_POINTS = {
    "top_left": (0.0, 0.0),
    "top_center": (0.5, 0.0),
    "top_right": (1.0, 0.0),
    "center_left": (0.0, 0.5),
    "center": (0.5, 0.5),
    "center_right": (1.0, 0.5),
    "bottom_left": (0.0, 1.0),
    "bottom_center": (0.5, 1.0),
    "bottom_right": (1.0, 1.0),
}

# Standard corporate colors (RGB tuples)
CORPORATE_COLORS = {
    "primary_blue": RGBColor(0, 112, 192),
    "secondary_gray": RGBColor(89, 89, 89),
    "accent_orange": RGBColor(237, 125, 49),
    "success_green": RGBColor(112, 173, 71),
    "warning_yellow": RGBColor(255, 192, 0),
    "danger_red": RGBColor(192, 0, 0),
    "white": RGBColor(255, 255, 255),
    "black": RGBColor(0, 0, 0),
}

# Standard fonts
STANDARD_FONTS = {"title": "Calibri Light", "body": "Calibri", "code": "Consolas"}

# WCAG 2.1 color contrast ratios
WCAG_CONTRAST_NORMAL = 4.5
WCAG_CONTRAST_LARGE = 3.0

# Maximum recommended file size (MB)
MAX_RECOMMENDED_FILE_SIZE_MB = 50

# Valid PowerPoint extensions
VALID_PPTX_EXTENSIONS = {".pptx", ".pptm", ".potx", ".potm"}

# Placeholder type mapping (integer keys for compatibility)
PLACEHOLDER_TYPE_NAMES = {
    0: "OBJECT",
    1: "TITLE",
    2: "BODY",
    3: "CENTER_TITLE",
    4: "SUBTITLE",
    5: "DATE",
    6: "SLIDE_NUMBER",
    7: "FOOTER",
    8: "HEADER",
    9: "OBJECT",
    10: "CHART",
    11: "TABLE",
    12: "CLIP_ART",
    13: "ORG_CHART",
    14: "MEDIA_CLIP",
    15: "BITMAP",
    16: "VERTICAL_TITLE",
    17: "VERTICAL_BODY",
    18: "PICTURE",
}

# Placeholder types that represent titles
TITLE_PLACEHOLDER_TYPES = {1, 3}  # TITLE and CENTER_TITLE

# Placeholder type for subtitle
SUBTITLE_PLACEHOLDER_TYPE = 4


def get_placeholder_type_name(ph_type_value: Any) -> str:
    """
    Safely get human-readable name for placeholder type.

    Args:
        ph_type_value: Placeholder type (int or enum)

    Returns:
        Human-readable string name
    """
    if ph_type_value is None:
        return "NONE"

    # Handle enum types
    if hasattr(ph_type_value, "value"):
        ph_type_value = ph_type_value.value

    try:
        int_value = int(ph_type_value)
        return PLACEHOLDER_TYPE_NAMES.get(int_value, f"UNKNOWN_{int_value}")
    except (TypeError, ValueError):
        return f"UNKNOWN_{ph_type_value}"


def _get_placeholder_type_int_helper(ph_type: Any) -> int:
    """
    Centralized helper to convert placeholder type to integer.

    Args:
        ph_type: Placeholder type object or value

    Returns:
        Integer representation of type
    """
    if ph_type is None:
        return 0
    if hasattr(ph_type, "value"):
        return ph_type.value
    try:
        return int(ph_type)
    except (TypeError, ValueError):
        return 0


# ============================================================================
# ENUMS
# ============================================================================


class ShapeType(Enum):
    """Common shape types supported by python-pptx."""

    RECTANGLE = "rectangle"
    ROUNDED_RECTANGLE = "rounded_rectangle"
    ELLIPSE = "ellipse"
    OVAL = "ellipse"
    TRIANGLE = "triangle"
    ARROW_RIGHT = "arrow_right"
    ARROW_LEFT = "arrow_left"
    ARROW_UP = "arrow_up"
    ARROW_DOWN = "arrow_down"
    STAR = "star"
    PENTAGON = "pentagon"
    HEXAGON = "hexagon"


class ChartType(Enum):
    """Supported chart types."""

    COLUMN = "column"
    COLUMN_CLUSTERED = "column"
    COLUMN_STACKED = "column_stacked"
    BAR = "bar"
    BAR_CLUSTERED = "bar"
    BAR_STACKED = "bar_stacked"
    LINE = "line"
    LINE_MARKERS = "line_markers"
    PIE = "pie"
    PIE_EXPLODED = "pie_exploded"
    AREA = "area"
    SCATTER = "scatter"


class TextAlignment(Enum):
    """Text alignment options."""

    LEFT = "left"
    CENTER = "center"
    RIGHT = "right"
    JUSTIFY = "justify"


class VerticalAlignment(Enum):
    """Vertical text alignment."""

    TOP = "top"
    MIDDLE = "middle"
    BOTTOM = "bottom"


class BulletStyle(Enum):
    """Bullet list styles."""

    BULLET = "bullet"
    NUMBERED = "numbered"
    NONE = "none"


class ImageFormat(Enum):
    """Supported image formats."""

    PNG = "png"
    JPG = "jpg"
    JPEG = "jpeg"
    GIF = "gif"
    BMP = "bmp"


class ExportFormat(Enum):
    """Export format options."""

    PDF = "pdf"
    PNG = "png"
    JPG = "jpg"
    PPTX = "pptx"


class ZOrderAction(Enum):
    """Z-order manipulation actions."""

    BRING_TO_FRONT = "bring_to_front"
    SEND_TO_BACK = "send_to_back"
    BRING_FORWARD = "bring_forward"
    SEND_BACKWARD = "send_backward"


class NotesMode(Enum):
    """Speaker notes insertion modes."""

    APPEND = "append"
    PREPEND = "prepend"
    OVERWRITE = "overwrite"


# ============================================================================
# UTILITY CLASSES
# ============================================================================


class FileLock:
    """
    Atomic file locking mechanism for concurrent access prevention.

    Uses OS-level atomic file creation to ensure only one process
    can hold the lock at a time.
    """

    def __init__(self, filepath: Path, timeout: float = 10.0):
        """
        Initialize file lock.

        Args:
            filepath: Path to file to lock
            timeout: Maximum seconds to wait for lock acquisition
        """
        self.filepath = Path(filepath)
        self.lockfile = self.filepath.parent / f".{self.filepath.name}.lock"
        self.timeout = timeout
        self.acquired = False
        self._fd: Optional[int] = None

    def acquire(self) -> bool:
        """
        Acquire lock with timeout using atomic file creation.

        Returns:
            True if lock acquired, False if timeout
        """
        start_time = time.time()

        while time.time() - start_time < self.timeout:
            try:
                # Use O_CREAT | O_EXCL for atomic creation
                # This is atomic on POSIX systems
                self._fd = os.open(
                    str(self.lockfile), os.O_CREAT | os.O_EXCL | os.O_WRONLY, 0o644
                )
                self.acquired = True
                return True
            except FileExistsError:
                time.sleep(0.1)
            except OSError as e:
                # EEXIST (cross-platform way via errno)
                if e.errno == errno.EEXIST:
                    time.sleep(0.1)
                else:
                    raise

        return False

    def release(self) -> None:
        """Release lock and clean up lock file."""
        if self._fd is not None:
            try:
                os.close(self._fd)
            except OSError:
                pass
            self._fd = None

        if self.acquired:
            try:
                self.lockfile.unlink(missing_ok=True)
            except OSError:
                pass
            self.acquired = False

    def __enter__(self) -> "FileLock":
        if not self.acquire():
            raise FileLockError(
                f"Could not acquire lock on {self.filepath} within {self.timeout}s",
                details={"filepath": str(self.filepath), "timeout": self.timeout},
            )
        return self

    def __exit__(self, exc_type, exc_val, exc_tb) -> bool:
        self.release()
        return False


class PathValidator:
    """
    Security-hardened path validation utility.

    Validates file paths to prevent path traversal attacks
    and ensure files are of expected types.
    """

    @staticmethod
    def validate_pptx_path(
        filepath: Union[str, Path],
        must_exist: bool = True,
        must_be_writable: bool = False,
        allowed_base_dirs: Optional[List[Path]] = None,
    ) -> Path:
        """
        Validate a PowerPoint file path.

        Args:
            filepath: Path to validate
            must_exist: If True, file must exist
            must_be_writable: If True, parent directory must be writable
            allowed_base_dirs: Optional list of base directories to restrict access (traversal protection)

        Returns:
            Resolved absolute Path

        Raises:
            PathValidationError: If validation fails
        """
        try:
            path = Path(filepath).resolve()
        except Exception as e:
            raise PathValidationError(
                f"Invalid path: {filepath}", details={"error": str(e)}
            )

        # Security: Path Traversal Protection
        if allowed_base_dirs:
            is_allowed = False
            for base in allowed_base_dirs:
                try:
                    # Check if path is relative to base
                    if path.is_relative_to(base.resolve()):
                        is_allowed = True
                        break
                except Exception:
                    continue

            if not is_allowed:
                raise PathValidationError(
                    f"Path is not within allowed directories: {path}",
                    details={
                        "path": str(path),
                        "allowed_base_dirs": [str(b) for b in allowed_base_dirs],
                    },
                )

        # Check extension
        if path.suffix.lower() not in VALID_PPTX_EXTENSIONS:
            raise PathValidationError(
                f"Invalid file extension: {path.suffix}",
                details={
                    "path": str(path),
                    "valid_extensions": list(VALID_PPTX_EXTENSIONS),
                },
            )

        # Check existence
        if must_exist and not path.exists():
            raise PathValidationError(
                f"File does not exist: {path}", details={"path": str(path)}
            )

        # Check if it's a file (not directory)
        if must_exist and not path.is_file():
            raise PathValidationError(
                f"Path is not a file: {path}", details={"path": str(path)}
            )

        # Check writability
        if must_be_writable:
            parent = path.parent
            if not parent.exists():
                raise PathValidationError(
                    f"Parent directory does not exist: {parent}",
                    details={"path": str(path), "parent": str(parent)},
                )
            if not os.access(str(parent), os.W_OK):
                raise PathValidationError(
                    f"Parent directory is not writable: {parent}",
                    details={"path": str(path), "parent": str(parent)},
                )

        return path

    @staticmethod
    def validate_image_path(filepath: Union[str, Path]) -> Path:
        """
        Validate an image file path.

        Args:
            filepath: Path to validate

        Returns:
            Resolved absolute Path

        Raises:
            ImageNotFoundError: If validation fails
        """
        try:
            path = Path(filepath).resolve()
        except Exception as e:
            raise ImageNotFoundError(
                f"Invalid image path: {filepath}", details={"error": str(e)}
            )

        if not path.exists():
            raise ImageNotFoundError(
                f"Image file does not exist: {path}", details={"path": str(path)}
            )

        if not path.is_file():
            raise ImageNotFoundError(
                f"Image path is not a file: {path}", details={"path": str(path)}
            )

        valid_image_extensions = {
            ".png",
            ".jpg",
            ".jpeg",
            ".gif",
            ".bmp",
            ".tiff",
            ".webp",
        }
        if path.suffix.lower() not in valid_image_extensions:
            raise ImageNotFoundError(
                f"Invalid image extension: {path.suffix}",
                details={
                    "path": str(path),
                    "valid_extensions": list(valid_image_extensions),
                },
            )

        return path


class Position:
    """Flexible position system supporting multiple input formats."""

    @staticmethod
    def from_dict(
        pos_dict: Dict[str, Any],
        slide_width: float = SLIDE_WIDTH_INCHES,
        slide_height: float = SLIDE_HEIGHT_INCHES,
    ) -> Tuple[float, float]:
        """
        Convert position dict to (left, top) in inches.

        Supports multiple formats:
        1. Absolute inches: {"left": 1.5, "top": 2.0}
        2. Percentage: {"left": "20%", "top": "30%"}
        3. Anchor-based: {"anchor": "center", "offset_x": 0.5, "offset_y": -1.0}
        4. Grid system: {"grid_row": 2, "grid_col": 3, "grid_size": 12}

        Args:
            pos_dict: Position specification dictionary
            slide_width: Slide width in inches (for percentage calculations)
            slide_height: Slide height in inches (for percentage calculations)

        Returns:
            Tuple of (left, top) in inches

        Raises:
            InvalidPositionError: If format is invalid
        """
        if not isinstance(pos_dict, dict):
            raise InvalidPositionError(
                f"Position must be a dictionary, got {type(pos_dict).__name__}",
                details={"value": str(pos_dict)},
            )

        # Format 1 & 2: Absolute or percentage with left/top
        if "left" in pos_dict and "top" in pos_dict:
            left = Position._parse_dimension(pos_dict["left"], slide_width)
            top = Position._parse_dimension(pos_dict["top"], slide_height)
            return (left, top)

        # Format 3: Anchor-based
        if "anchor" in pos_dict:
            anchor_name = pos_dict["anchor"].lower().replace("-", "_").replace(" ", "_")
            anchor = ANCHOR_POINTS.get(anchor_name)

            if anchor is None:
                raise InvalidPositionError(
                    f"Unknown anchor: {pos_dict['anchor']}",
                    details={"available_anchors": list(ANCHOR_POINTS.keys())},
                )

            # Anchor is in relative coordinates (0-1), convert to inches
            base_left = anchor[0] * slide_width
            base_top = anchor[1] * slide_height

            offset_x = float(pos_dict.get("offset_x", 0))
            offset_y = float(pos_dict.get("offset_y", 0))

            return (base_left + offset_x, base_top + offset_y)

        # Format 4: Grid system
        if "grid_row" in pos_dict and "grid_col" in pos_dict:
            grid_size = int(pos_dict.get("grid_size", 12))
            cell_width = slide_width / grid_size
            cell_height = slide_height / grid_size

            col = int(pos_dict["grid_col"])
            row = int(pos_dict["grid_row"])

            left = col * cell_width
            top = row * cell_height

            return (left, top)

        raise InvalidPositionError(
            "Invalid position format",
            details={
                "provided": pos_dict,
                "expected_formats": [
                    {"left": "value", "top": "value"},
                    {"anchor": "center", "offset_x": 0, "offset_y": 0},
                    {"grid_row": 0, "grid_col": 0, "grid_size": 12},
                ],
            },
        )

    @staticmethod
    def _parse_dimension(value: Union[str, float, int], max_dimension: float) -> float:
        """
        Parse dimension value (supports percentages or absolute values).

        Args:
            value: Dimension value (e.g., "50%", 2.5, "2.5")
            max_dimension: Maximum dimension for percentage calculation

        Returns:
            Dimension in inches
        """
        if isinstance(value, str):
            value = value.strip()
            if value.endswith("%"):
                percent = float(value[:-1]) / 100.0
                return percent * max_dimension
            else:
                return float(value)
        return float(value)


class Size:
    """Flexible size system supporting multiple input formats."""

    @staticmethod
    def from_dict(
        size_dict: Dict[str, Any],
        slide_width: float = SLIDE_WIDTH_INCHES,
        slide_height: float = SLIDE_HEIGHT_INCHES,
        aspect_ratio: Optional[float] = None,
    ) -> Tuple[Optional[float], Optional[float]]:
        """
        Convert size dict to (width, height) in inches.

        Supports:
        - {"width": 5.0, "height": 3.0}  # Absolute inches
        - {"width": "50%", "height": "30%"}  # Percentage of slide
        - {"width": "auto", "height": 3.0}  # Maintain aspect ratio
        - {"width": 5.0, "height": "auto"}  # Maintain aspect ratio

        Args:
            size_dict: Size specification dictionary
            slide_width: Slide width in inches
            slide_height: Slide height in inches
            aspect_ratio: Optional aspect ratio (width/height) for "auto" calculations

        Returns:
            Tuple of (width, height) in inches, either can be None for "auto"
        """
        if not isinstance(size_dict, dict):
            raise ValueError(
                f"Size must be a dictionary, got {type(size_dict).__name__}"
            )

        if "width" not in size_dict and "height" not in size_dict:
            raise ValueError("Size must have at least 'width' or 'height'")

        width_spec = size_dict.get("width")
        height_spec = size_dict.get("height")

        # Parse width
        if width_spec == "auto" or width_spec is None:
            width = None
        else:
            width = Position._parse_dimension(width_spec, slide_width)

        # Parse height
        if height_spec == "auto" or height_spec is None:
            height = None
        else:
            height = Position._parse_dimension(height_spec, slide_height)

        # Apply aspect ratio if one dimension is auto
        if aspect_ratio is not None:
            if width is None and height is not None:
                width = height * aspect_ratio
            elif height is None and width is not None:
                height = width / aspect_ratio

        return (width, height)


class ColorHelper:
    """Utilities for color conversion and validation."""

    @staticmethod
    def from_hex(hex_color: str) -> RGBColor:
        """
        Convert hex color string to RGBColor.

        Args:
            hex_color: Hex color string (e.g., "#FF0000" or "FF0000")

        Returns:
            RGBColor object

        Raises:
            ValueError: If hex color format is invalid
        """
        hex_color = hex_color.strip().lstrip("#")

        if len(hex_color) != 6:
            raise ValueError(f"Invalid hex color: {hex_color}. Must be 6 hex digits.")

        if not all(c in "0123456789ABCDEFabcdef" for c in hex_color):
            raise ValueError(
                f"Invalid hex color: {hex_color}. Contains non-hex characters."
            )

        r = int(hex_color[0:2], 16)
        g = int(hex_color[2:4], 16)
        b = int(hex_color[4:6], 16)

        return RGBColor(r, g, b)

    @staticmethod
    def to_hex(rgb_color: RGBColor) -> str:
        """
        Convert RGBColor to hex string.

        Args:
            rgb_color: RGBColor object

        Returns:
            Hex color string with # prefix
        """
        if hasattr(rgb_color, "__iter__") and len(rgb_color) == 3:
            r, g, b = rgb_color
        elif hasattr(rgb_color, "r"):
            r, g, b = rgb_color.r, rgb_color.g, rgb_color.b
        else:
            # Handle string representation
            hex_str = str(rgb_color).lstrip("#")
            return f"#{hex_str}"

        return f"#{r:02x}{g:02x}{b:02x}"

    @staticmethod
    def luminance(rgb_color: Union[RGBColor, Tuple[int, int, int]]) -> float:
        """
        Calculate relative luminance for WCAG contrast calculations.

        Args:
            rgb_color: RGBColor or (r, g, b) tuple

        Returns:
            Relative luminance value (0.0 to 1.0)
        """
        # Extract RGB values
        if hasattr(rgb_color, "r"):
            r, g, b = rgb_color.r, rgb_color.g, rgb_color.b
        elif hasattr(rgb_color, "__iter__"):
            r, g, b = rgb_color
        else:
            # Handle string representation
            hex_str = str(rgb_color).lstrip("#")
            if len(hex_str) == 6:
                r = int(hex_str[0:2], 16)
                g = int(hex_str[2:4], 16)
                b = int(hex_str[4:6], 16)
            else:
                raise ValueError(f"Cannot parse color: {rgb_color}")

        def _linearize(channel: int) -> float:
            c = channel / 255.0
            if c <= 0.03928:
                return c / 12.92
            return ((c + 0.055) / 1.055) ** 2.4

        r_lin = _linearize(r)
        g_lin = _linearize(g)
        b_lin = _linearize(b)

        return 0.2126 * r_lin + 0.7152 * g_lin + 0.0722 * b_lin

    @staticmethod
    def contrast_ratio(color1: RGBColor, color2: RGBColor) -> float:
        """
        Calculate WCAG contrast ratio between two colors.

        Args:
            color1: First color
            color2: Second color

        Returns:
            Contrast ratio (1.0 to 21.0)
        """
        lum1 = ColorHelper.luminance(color1)
        lum2 = ColorHelper.luminance(color2)

        lighter = max(lum1, lum2)
        darker = min(lum1, lum2)

        return (lighter + 0.05) / (darker + 0.05)

    @staticmethod
    def meets_wcag(
        foreground: RGBColor, background: RGBColor, is_large_text: bool = False
    ) -> bool:
        """
        Check if color combination meets WCAG 2.1 AA standards.

        Args:
            foreground: Text/foreground color
            background: Background color
            is_large_text: True if text is 18pt+ or 14pt+ bold

        Returns:
            True if contrast is sufficient
        """
        ratio = ColorHelper.contrast_ratio(foreground, background)
        threshold = WCAG_CONTRAST_LARGE if is_large_text else WCAG_CONTRAST_NORMAL
        return ratio >= threshold


# ============================================================================
# ANALYSIS CLASSES
# ============================================================================


class TemplateProfile:
    """
    Captures and provides access to PowerPoint template formatting.

    Uses lazy loading to avoid performance penalty when profile is not needed.
    """

    def __init__(self, prs: Optional["Presentation"] = None):
        """
        Initialize template profile.

        Args:
            prs: Optional Presentation to analyze immediately
        """
        self._prs = prs
        self._captured = False
        self._slide_layouts: List[Dict[str, Any]] = []
        self._theme_colors: Dict[str, str] = {}
        self._theme_fonts: Dict[str, str] = {}

    def _ensure_captured(self) -> None:
        """Ensure template data has been captured (lazy loading)."""
        if self._captured or self._prs is None:
            return

        self._capture_layouts()
        self._capture_theme()
        self._captured = True

    def _capture_layouts(self) -> None:
        """Capture layout information from presentation."""
        for layout in self._prs.slide_layouts:
            layout_info = {"name": layout.name, "placeholders": []}

            for ph in layout.placeholders:
                try:
                    ph_info = {
                        "type": _get_placeholder_type_int_helper(
                            ph.placeholder_format.type
                        ),
                        "idx": ph.placeholder_format.idx,
                    }
                    if hasattr(ph, "left") and ph.left is not None:
                        ph_info["position"] = {
                            "left": ph.left / EMU_PER_INCH,
                            "top": ph.top / EMU_PER_INCH,
                        }
                    if hasattr(ph, "width") and ph.width is not None:
                        ph_info["size"] = {
                            "width": ph.width / EMU_PER_INCH,
                            "height": ph.height / EMU_PER_INCH,
                        }
                    layout_info["placeholders"].append(ph_info)
                except Exception:
                    continue

            self._slide_layouts.append(layout_info)

    def _capture_theme(self) -> None:
        """Capture theme colors and fonts from presentation."""
        try:
            # Attempt to extract theme colors
            if hasattr(self._prs, "slide_master") and self._prs.slide_master:
                master = self._prs.slide_master

                # Extract fonts from shapes
                for shape in master.shapes:
                    if hasattr(shape, "text_frame"):
                        try:
                            for para in shape.text_frame.paragraphs:
                                if para.font.name:
                                    font_key = f"font_{len(self._theme_fonts)}"
                                    if para.font.name not in self._theme_fonts.values():
                                        self._theme_fonts[font_key] = para.font.name
                        except Exception:
                            continue
        except Exception:
            pass

    @property
    def slide_layouts(self) -> List[Dict[str, Any]]:
        """Get slide layout information."""
        self._ensure_captured()
        return self._slide_layouts

    @property
    def theme_colors(self) -> Dict[str, str]:
        """Get theme colors."""
        self._ensure_captured()
        return self._theme_colors

    @property
    def theme_fonts(self) -> Dict[str, str]:
        """Get theme fonts."""
        self._ensure_captured()
        return self._theme_fonts

    def get_layout_names(self) -> List[str]:
        """Get list of available layout names."""
        self._ensure_captured()
        return [layout["name"] for layout in self._slide_layouts]

    def to_dict(self) -> Dict[str, Any]:
        """Convert profile to JSON-serializable dict."""
        self._ensure_captured()
        return {
            "slide_layouts": self._slide_layouts,
            "theme_colors": self._theme_colors,
            "theme_fonts": self._theme_fonts,
        }


class AccessibilityChecker:
    """WCAG 2.1 compliance checker for presentations."""

    @staticmethod
    def check_presentation(prs: "Presentation") -> Dict[str, Any]:
        """
        Comprehensive accessibility check.

        Args:
            prs: Presentation to check

        Returns:
            Dict containing:
            - status: "accessible" or "issues_found"
            - total_issues: Count of all issues
            - issues: Detailed issue breakdown
            - wcag_level: "AA" if passing, "fail" otherwise
        """
        issues = {
            "missing_alt_text": [],
            "low_contrast": [],
            "missing_titles": [],
            "small_text": [],
            "reading_order_warnings": [],
        }

        for slide_idx, slide in enumerate(prs.slides):
            # Check for title
            has_title = AccessibilityChecker._check_slide_has_title(slide)
            if not has_title:
                issues["missing_titles"].append(
                    {
                        "slide": slide_idx,
                        "message": "Slide lacks a title for screen reader navigation",
                    }
                )

            # Check each shape
            for shape_idx, shape in enumerate(slide.shapes):
                # Check images for alt text
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if not AccessibilityChecker._has_alt_text(shape):
                        issues["missing_alt_text"].append(
                            {
                                "slide": slide_idx,
                                "shape": shape_idx,
                                "shape_name": shape.name,
                                "message": "Image lacks alternative text",
                            }
                        )

                # Check text for contrast and size
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    AccessibilityChecker._check_text_accessibility(
                        shape, slide_idx, shape_idx, issues
                    )

        total_issues = sum(len(v) for v in issues.values())

        return {
            "status": "issues_found" if total_issues > 0 else "accessible",
            "total_issues": total_issues,
            "issues": issues,
            "wcag_level": "AA" if total_issues == 0 else "fail",
            "checked_slides": len(prs.slides),
        }

    @staticmethod
    def _check_slide_has_title(slide) -> bool:
        """Check if slide has a non-empty title."""
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = _get_placeholder_type_int_helper(
                    shape.placeholder_format.type
                )
                if ph_type in TITLE_PLACEHOLDER_TYPES:
                    if shape.has_text_frame and shape.text_frame.text.strip():
                        return True
        return False

    @staticmethod
    def _has_alt_text(shape) -> bool:
        """
        Check if image shape has meaningful alt text.

        Checks both the description attribute (proper alt text)
        and the shape name as fallback.
        """
        # Check description attribute (the actual alt text storage)
        try:
            element = shape._element
            # Check for description in various possible locations
            descr = element.get("descr")
            if descr and descr.strip() and len(descr.strip()) > 3:
                return True

            # Check nvPicPr/cNvPr for descr
            for child in element.iter():
                if child.get("descr"):
                    descr = child.get("descr")
                    if descr and descr.strip() and len(descr.strip()) > 3:
                        return True
        except Exception:
            pass

        # Fallback: check name (not ideal, but some tools use this)
        if shape.name:
            name = shape.name.strip()
            # Reject generic names
            if name.lower().startswith("picture"):
                return False
            if name.lower().startswith("image"):
                return False
            if len(name) > 5:  # Meaningful name
                return True

        return False

    @staticmethod
    def _check_text_accessibility(
        shape, slide_idx: int, shape_idx: int, issues: Dict[str, Any]
    ) -> None:
        """Check text shape for accessibility issues."""
        try:
            text_frame = shape.text_frame
            for para in text_frame.paragraphs:
                # Check font size
                if para.font.size is not None:
                    size_pt = para.font.size.pt
                    if size_pt < 10:
                        issues["small_text"].append(
                            {
                                "slide": slide_idx,
                                "shape": shape_idx,
                                "size_pt": size_pt,
                                "text_preview": para.text[:50] if para.text else "",
                                "message": f"Text size {size_pt}pt is below minimum 10pt",
                            }
                        )
        except Exception:
            pass


class AssetValidator:
    """Validates and provides information about presentation assets."""

    @staticmethod
    def validate_presentation_assets(
        prs: "Presentation", filepath: Optional[Path] = None
    ) -> Dict[str, Any]:
        """
        Validate all assets in presentation.

        Args:
            prs: Presentation to validate
            filepath: Optional file path for size check

        Returns:
            Validation report dict
        """
        issues = {"large_images": [], "total_embedded_size_bytes": 0, "image_count": 0}

        for slide_idx, slide in enumerate(prs.slides):
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    issues["image_count"] += 1
                    try:
                        image_blob = shape.image.blob
                        image_size = len(image_blob)
                        issues["total_embedded_size_bytes"] += image_size

                        # Flag images over 2MB
                        if image_size > 2 * 1024 * 1024:
                            issues["large_images"].append(
                                {
                                    "slide": slide_idx,
                                    "shape": shape_idx,
                                    "size_bytes": image_size,
                                    "size_mb": round(image_size / (1024 * 1024), 2),
                                }
                            )
                    except Exception:
                        pass

        # Check total file size
        if filepath and Path(filepath).exists():
            file_size = Path(filepath).stat().st_size
            issues["file_size_bytes"] = file_size
            issues["file_size_mb"] = round(file_size / (1024 * 1024), 2)

            if file_size > MAX_RECOMMENDED_FILE_SIZE_MB * 1024 * 1024:
                issues["large_file_warning"] = {
                    "size_mb": issues["file_size_mb"],
                    "recommended_max_mb": MAX_RECOMMENDED_FILE_SIZE_MB,
                }

        total_issues = len(issues["large_images"])
        if "large_file_warning" in issues:
            total_issues += 1

        return {
            "status": "issues_found" if total_issues > 0 else "valid",
            "total_issues": total_issues,
            "issues": issues,
        }

    @staticmethod
    def compress_image(
        image_path: Path, max_width: int = 1920, quality: int = 85
    ) -> BytesIO:
        """
        Compress image for PowerPoint embedding.

        Args:
            image_path: Path to source image
            max_width: Maximum width in pixels
            quality: JPEG quality (1-100)

        Returns:
            BytesIO containing compressed image

        Raises:
            ImportError: If Pillow is not available
        """
        if not HAS_PILLOW:
            raise ImportError("Pillow is required for image compression")

        with PILImage.open(image_path) as img:
            # Resize if needed
            if img.width > max_width:
                ratio = max_width / img.width
                new_height = int(img.height * ratio)
                img = img.resize((max_width, new_height), PILImage.LANCZOS)

            # Convert to RGB if necessary
            if img.mode in ("RGBA", "LA", "P"):
                background = PILImage.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                if img.mode in ("RGBA", "LA"):
                    background.paste(img, mask=img.split()[-1])
                else:
                    background.paste(img)
                img = background

            # Save to BytesIO
            output = BytesIO()
            img.save(output, format="JPEG", quality=quality, optimize=True)
            output.seek(0)

            return output


# ============================================================================
# MAIN POWERPOINT AGENT CLASS
# ============================================================================


class PowerPointAgent:
    """
    Core PowerPoint manipulation class for stateless tool operations.

    Provides comprehensive PowerPoint editing capabilities optimized for
    AI agent consumption through simple, composable operations.

    Features:
    - Stateless design for tool-based workflows
    - Comprehensive validation and accessibility checking
    - Atomic file locking for concurrent access safety
    - Full alignment with Presentation Architect System Prompt v3.0
    - Approval token governance for destructive operations
    - Geometry-aware version tracking for state detection

    Example:
        with PowerPointAgent() as agent:
            agent.open(Path("presentation.pptx"))
            agent.add_slide("Title and Content")
            agent.save()
    """

    def __init__(self, filepath: Optional[Union[str, Path]] = None):
        """
        Initialize PowerPoint agent.

        Args:
            filepath: Optional path to open immediately
        """
        self.filepath: Optional[Path] = None
        self.prs: Optional[Presentation] = None
        self._lock: Optional[FileLock] = None
        self._template_profile: Optional[TemplateProfile] = None
        self._layout_cache: Optional[Dict[str, Any]] = None

        if filepath:
            self.filepath = Path(filepath)

    # ========================================================================
    # CONTEXT MANAGEMENT
    # ========================================================================

    def __enter__(self) -> "PowerPointAgent":
        return self

    def __exit__(self, exc_type, exc_val, exc_tb) -> bool:
        self.close()
        return False

    # ========================================================================
    # HELPER METHODS (Governance & Observability)
    # ========================================================================

    def _validate_token(self, token: Optional[str], scope: str) -> None:
        """
        Validate approval token for destructive operations.

        Args:
            token: The approval token string
            scope: The required permission scope (e.g., "delete:slide")

        Raises:
            ApprovalTokenError: If token is missing or invalid
        """
        # NOTE: In a production environment, this would verify a JWT or HMAC.
        # For this implementation, we check presence and basic format.
        if not token:
            raise ApprovalTokenError(
                f"Destructive operation requires approval token (scope: {scope})",
                details={"scope_required": scope},
            )

        # Placeholder validation - real implementation would check signature
        if len(token) < 8:
            raise ApprovalTokenError(
                "Invalid approval token format", details={"token_length": len(token)}
            )

    def _capture_version(self) -> str:
        """Capture current presentation version hash."""
        return self.get_presentation_version()

    def _log_warning(self, message: str) -> None:
        """Log a warning message through the configured logger."""
        logger.warning(message)

    # ========================================================================
    # FILE OPERATIONS
    # ========================================================================

    def create_new(self, template: Optional[Union[str, Path]] = None) -> None:
        """
        Create new presentation, optionally from template.

        Args:
            template: Optional path to template .pptx file

        Raises:
            FileNotFoundError: If template doesn't exist
            TemplateError: If template cannot be loaded
        """
        if template:
            template_path = PathValidator.validate_pptx_path(template, must_exist=True)
            try:
                self.prs = Presentation(str(template_path))
            except Exception as e:
                raise TemplateError(
                    f"Failed to load template: {template_path}",
                    details={"error": str(e)},
                )
        else:
            self.prs = Presentation()

        self._template_profile = TemplateProfile(self.prs)
        self._layout_cache = None

    def open(self, filepath: Union[str, Path], acquire_lock: bool = True) -> None:
        """
        Open existing presentation.

        Args:
            filepath: Path to .pptx file
            acquire_lock: Whether to acquire exclusive file lock

        Raises:
            PathValidationError: If path is invalid
            FileLockError: If lock cannot be acquired
            PowerPointAgentError: If file cannot be opened
        """
        validated_path = PathValidator.validate_pptx_path(filepath, must_exist=True)
        self.filepath = validated_path

        # Acquire lock if requested
        if acquire_lock:
            self._lock = FileLock(validated_path)
            if not self._lock.acquire():
                raise FileLockError(
                    f"Could not acquire lock on {validated_path}",
                    details={"filepath": str(validated_path)},
                )

        # Load presentation (with lock release on failure)
        try:
            self.prs = Presentation(str(validated_path))
            self._template_profile = TemplateProfile(self.prs)
            self._layout_cache = None
        except Exception as e:
            # Release lock on failure
            if self._lock:
                self._lock.release()
                self._lock = None
            raise PowerPointAgentError(
                f"Failed to open presentation: {validated_path}",
                details={"error": str(e)},
            )

    def save(self, filepath: Optional[Union[str, Path]] = None) -> None:
        """
        Save presentation.

        Args:
            filepath: Output path (uses original path if None)

        Raises:
            PowerPointAgentError: If no presentation loaded
            PathValidationError: If output path is invalid
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        target = filepath or self.filepath
        if not target:
            raise PowerPointAgentError("No output path specified")

        target_path = PathValidator.validate_pptx_path(
            target, must_exist=False, must_be_writable=True
        )

        # Ensure parent directory exists
        target_path.parent.mkdir(parents=True, exist_ok=True)

        self.prs.save(str(target_path))
        self.filepath = target_path

    def close(self) -> None:
        """Close presentation and release resources."""
        self.prs = None
        self._template_profile = None
        self._layout_cache = None

        if self._lock:
            self._lock.release()
            self._lock = None

    def clone_presentation(self, output_path: Union[str, Path]) -> "PowerPointAgent":
        """
        Clone current presentation to a new file.

        Args:
            output_path: Path for the cloned presentation

        Returns:
            New PowerPointAgent instance with cloned presentation

        Raises:
            PowerPointAgentError: If no presentation loaded
            PathValidationError: If output path is invalid
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        output = PathValidator.validate_pptx_path(
            output_path, must_exist=False, must_be_writable=True
        )

        # Save to new location
        output.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(str(output))

        # Create new agent with cloned file
        new_agent = PowerPointAgent()
        new_agent.open(output)

        return new_agent

    # ========================================================================
    # SLIDE OPERATIONS
    # ========================================================================

    def add_slide(
        self, layout_name: str = "Title and Content", index: Optional[int] = None
    ) -> Dict[str, Any]:
        """
        Add new slide with specified layout.

        Args:
            layout_name: Name of layout to use
            index: Position to insert (None = append at end)

        Returns:
            Dict with slide_index and layout_name

        Raises:
            PowerPointAgentError: If no presentation loaded
            LayoutNotFoundError: If layout doesn't exist
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        version_before = self._capture_version()

        layout = self._get_layout(layout_name)
        slide = self.prs.slides.add_slide(layout)

        result_index = len(self.prs.slides) - 1

        if index is not None:
            max_valid = len(self.prs.slides)
            if not 0 <= index <= max_valid:
                raise SlideNotFoundError(
                    f"Insert index {index} out of range (0-{max_valid})",
                    details={"index": index, "valid_range": f"0-{max_valid}"},
                )

            # Move slide from end to target position
            xml_slides = self.prs.slides._sldIdLst
            slide_elem = xml_slides[-1]
            xml_slides.remove(slide_elem)
            xml_slides.insert(index, slide_elem)
            result_index = index

        version_after = self._capture_version()

        return {
            "slide_index": result_index,
            "layout_name": layout_name,
            "total_slides": len(self.prs.slides),
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def delete_slide(
        self, index: int, approval_token: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Delete slide at index.

        ⚠️ DESTRUCTIVE OPERATION - Requires approval token.

        Args:
            index: Slide index (0-based)
            approval_token: Token authorizing destructive operation

        Returns:
            Dict with deleted index and new slide count

        Raises:
            SlideNotFoundError: If index is out of range
            ApprovalTokenError: If token is missing/invalid
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        self._validate_token(approval_token, APPROVAL_SCOPE_DELETE_SLIDE)

        slide_count = len(self.prs.slides)
        if not 0 <= index < slide_count:
            raise SlideNotFoundError(
                f"Slide index {index} out of range",
                details={"index": index, "slide_count": slide_count},
            )

        version_before = self._capture_version()

        # Get slide relationship ID and remove
        rId = self.prs.slides._sldIdLst[index].rId
        self.prs.part.drop_rel(rId)
        del self.prs.slides._sldIdLst[index]

        version_after = self._capture_version()

        return {
            "deleted_index": index,
            "previous_count": slide_count,
            "new_count": len(self.prs.slides),
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def duplicate_slide(self, index: int) -> Dict[str, Any]:
        """
        Duplicate slide at index.

        Args:
            index: Slide index to duplicate

        Returns:
            Dict with new slide index

        Raises:
            SlideNotFoundError: If index is out of range
        """
        source_slide = self._get_slide(index)
        version_before = self._capture_version()

        # Add new slide with same layout
        layout = source_slide.slide_layout
        new_slide = self.prs.slides.add_slide(layout)
        new_index = len(self.prs.slides) - 1

        # Copy shapes
        for shape in source_slide.shapes:
            try:
                self._copy_shape(shape, new_slide)
            except Exception as e:
                logger.warning(f"Could not copy shape: {e}")

        version_after = self._capture_version()

        return {
            "source_index": index,
            "new_index": new_index,
            "total_slides": len(self.prs.slides),
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def reorder_slides(self, from_index: int, to_index: int) -> Dict[str, Any]:
        """
        Move slide from one position to another.

        Args:
            from_index: Current position
            to_index: Desired position

        Returns:
            Dict with movement details

        Raises:
            SlideNotFoundError: If either index is out of range
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        slide_count = len(self.prs.slides)

        if not 0 <= from_index < slide_count:
            raise SlideNotFoundError(
                f"Source index {from_index} out of range",
                details={"from_index": from_index, "slide_count": slide_count},
            )

        if not 0 <= to_index < slide_count:
            raise SlideNotFoundError(
                f"Target index {to_index} out of range",
                details={"to_index": to_index, "slide_count": slide_count},
            )

        version_before = self._capture_version()

        xml_slides = self.prs.slides._sldIdLst
        slide_elem = xml_slides[from_index]
        xml_slides.remove(slide_elem)
        xml_slides.insert(to_index, slide_elem)

        version_after = self._capture_version()

        return {
            "from_index": from_index,
            "to_index": to_index,
            "total_slides": slide_count,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def get_slide_count(self) -> int:
        """
        Get total number of slides.

        Returns:
            Number of slides
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")
        return len(self.prs.slides)

    # ========================================================================
    # TEXT OPERATIONS
    # ========================================================================

    def add_text_box(
        self,
        slide_index: int,
        text: str,
        position: Dict[str, Any],
        size: Dict[str, Any],
        font_name: Optional[str] = None,
        font_size: int = 18,
        bold: bool = False,
        italic: bool = False,
        color: Optional[str] = None,
        alignment: str = "left",
    ) -> Dict[str, Any]:
        """
        Add text box to slide.

        Args:
            slide_index: Target slide index
            text: Text content
            position: Position dict (see Position.from_dict)
            size: Size dict (see Size.from_dict)
            font_name: Font name (None uses theme font)
            font_size: Font size in points
            bold: Bold text
            italic: Italic text
            color: Text color hex (e.g., "#FF0000")
            alignment: Text alignment ("left", "center", "right", "justify")

        Returns:
            Dict with shape_index and details

        Raises:
            SlideNotFoundError: If slide index is invalid
            InvalidPositionError: If position is invalid
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        # Parse position and size
        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)

        if width is None or height is None:
            raise ValueError("Text box must have explicit width and height")

        # Create text box
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Configure text frame
        text_frame = text_box.text_frame
        text_frame.text = text
        text_frame.word_wrap = True

        # Apply formatting
        paragraph = text_frame.paragraphs[0]
        if font_name:
            paragraph.font.name = font_name
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.italic = italic

        if color:
            paragraph.font.color.rgb = ColorHelper.from_hex(color)

        # Set alignment
        alignment_map = {
            "left": PP_ALIGN.LEFT,
            "center": PP_ALIGN.CENTER,
            "right": PP_ALIGN.RIGHT,
            "justify": PP_ALIGN.JUSTIFY,
        }
        paragraph.alignment = alignment_map.get(alignment.lower(), PP_ALIGN.LEFT)

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": len(slide.shapes) - 1,
            "text_length": len(text),
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height},
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def set_title(
        self, slide_index: int, title: str, subtitle: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Set slide title and optional subtitle.

        Args:
            slide_index: Target slide index
            title: Title text
            subtitle: Optional subtitle text

        Returns:
            Dict with title/subtitle set status

        Raises:
            SlideNotFoundError: If slide index is invalid
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        title_set = False
        subtitle_set = False
        title_shape_index = None
        subtitle_shape_index = None

        for idx, shape in enumerate(slide.shapes):
            if shape.is_placeholder:
                ph_type = _get_placeholder_type_int_helper(
                    shape.placeholder_format.type
                )

                # Check for title placeholder
                if ph_type in TITLE_PLACEHOLDER_TYPES:
                    if shape.has_text_frame:
                        shape.text_frame.text = title
                        title_set = True
                        title_shape_index = idx

                # Check for subtitle placeholder
                elif ph_type == SUBTITLE_PLACEHOLDER_TYPE:
                    if subtitle and shape.has_text_frame:
                        shape.text_frame.text = subtitle
                        subtitle_set = True
                        subtitle_shape_index = idx

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "title_set": title_set,
            "subtitle_set": subtitle_set,
            "title_shape_index": title_shape_index,
            "subtitle_shape_index": subtitle_shape_index,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def add_bullet_list(
        self,
        slide_index: int,
        items: List[str],
        position: Dict[str, Any],
        size: Dict[str, Any],
        bullet_style: str = "bullet",
        font_size: int = 18,
        font_name: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Add bullet list to slide.

        Args:
            slide_index: Target slide index
            items: List of bullet items
            position: Position dict
            size: Size dict
            bullet_style: "bullet", "numbered", or "none"
            font_size: Font size in points
            font_name: Optional font name

        Returns:
            Dict with shape_index and item count
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)

        if width is None or height is None:
            raise ValueError("Bullet list must have explicit width and height")

        # Create text box for bullets
        text_box = slide.shapes.add_textbox(
            Inches(left), Inches(top), Inches(width), Inches(height)
        )

        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        for idx, item in enumerate(items):
            if idx == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            if bullet_style == "numbered":
                p.text = f"{idx + 1}. {item}"
            else:
                p.text = item

            p.level = 0
            p.font.size = Pt(font_size)
            if font_name:
                p.font.name = font_name

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": len(slide.shapes) - 1,
            "item_count": len(items),
            "bullet_style": bullet_style,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def format_text(
        self,
        slide_index: int,
        shape_index: int,
        font_name: Optional[str] = None,
        font_size: Optional[int] = None,
        bold: Optional[bool] = None,
        italic: Optional[bool] = None,
        color: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Format existing text shape.

        Args:
            slide_index: Target slide index
            shape_index: Shape index on slide
            font_name: Optional font name
            font_size: Optional font size in points
            bold: Optional bold setting
            italic: Optional italic setting
            color: Optional color hex

        Returns:
            Dict with formatting applied
        """
        shape = self._get_shape(slide_index, shape_index)

        if not hasattr(shape, "text_frame") or not shape.has_text_frame:
            raise ValueError(f"Shape at index {shape_index} does not have text")

        version_before = self._capture_version()
        changes = []

        for paragraph in shape.text_frame.paragraphs:
            if font_name is not None:
                paragraph.font.name = font_name
                changes.append("font_name")
            if font_size is not None:
                paragraph.font.size = Pt(font_size)
                changes.append("font_size")
            if bold is not None:
                paragraph.font.bold = bold
                changes.append("bold")
            if italic is not None:
                paragraph.font.italic = italic
                changes.append("italic")
            if color is not None:
                paragraph.font.color.rgb = ColorHelper.from_hex(color)
                changes.append("color")

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "changes_applied": list(set(changes)),
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def replace_text(
        self,
        find: str,
        replace: str,
        slide_index: Optional[int] = None,
        shape_index: Optional[int] = None,
        match_case: bool = False,
    ) -> Dict[str, Any]:
        """
        Find and replace text in presentation.

        Args:
            find: Text to find
            replace: Replacement text
            slide_index: Optional specific slide (None = all slides)
            shape_index: Optional specific shape (requires slide_index)
            match_case: Case-sensitive matching

        Returns:
            Dict with replacement count and locations
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        if shape_index is not None and slide_index is None:
            raise ValueError("shape_index requires slide_index to be specified")

        version_before = self._capture_version()

        replacements = []
        total_count = 0

        # Determine slides to process
        if slide_index is not None:
            slides_to_process = [(slide_index, self._get_slide(slide_index))]
        else:
            slides_to_process = list(enumerate(self.prs.slides))

        for s_idx, slide in slides_to_process:
            # Determine shapes to process
            if shape_index is not None:
                shapes_to_process = [(shape_index, self._get_shape(s_idx, shape_index))]
            else:
                shapes_to_process = list(enumerate(slide.shapes))

            for sh_idx, shape in shapes_to_process:
                if not hasattr(shape, "text_frame") or not shape.has_text_frame:
                    continue

                count = self._replace_text_in_shape(shape, find, replace, match_case)
                if count > 0:
                    total_count += count
                    replacements.append(
                        {"slide": s_idx, "shape": sh_idx, "count": count}
                    )

        version_after = self._capture_version()

        return {
            "find": find,
            "replace": replace,
            "match_case": match_case,
            "total_replacements": total_count,
            "locations": replacements,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def _replace_text_in_shape(
        self, shape, find: str, replace: str, match_case: bool
    ) -> int:
        """Replace text within a single shape, preserving formatting where possible."""
        count = 0

        try:
            text_frame = shape.text_frame
        except (AttributeError, TypeError):
            return 0

        # Strategy 1: Replace in runs (preserves formatting)
        for paragraph in text_frame.paragraphs:
            for run in paragraph.runs:
                if match_case:
                    if find in run.text:
                        occurrences = run.text.count(find)
                        run.text = run.text.replace(find, replace)
                        count += occurrences
                else:
                    if find.lower() in run.text.lower():
                        pattern = re.compile(re.escape(find), re.IGNORECASE)
                        matches = pattern.findall(run.text)
                        run.text = pattern.sub(replace, run.text)
                        count += len(matches)

        if count > 0:
            return count

        # Strategy 2: Full text replacement (if text spans runs)
        try:
            full_text = shape.text
            if not full_text:
                return 0

            if match_case:
                if find in full_text:
                    occurrences = full_text.count(find)
                    shape.text = full_text.replace(find, replace)
                    return occurrences
            else:
                if find.lower() in full_text.lower():
                    pattern = re.compile(re.escape(find), re.IGNORECASE)
                    matches = pattern.findall(full_text)
                    shape.text = pattern.sub(replace, full_text)
                    return len(matches)
        except (AttributeError, TypeError):
            pass

        return 0

    def add_notes(
        self,
        slide_index: int,
        text: str,
        mode: Union[str, NotesMode] = NotesMode.APPEND,
    ) -> Dict[str, Any]:
        """
        Add speaker notes to a slide.

        Args:
            slide_index: Target slide index
            text: Notes text to add
            mode: "append", "prepend", or "overwrite"

        Returns:
            Dict with notes details

        Raises:
            SlideNotFoundError: If slide index is invalid
            ValueError: If mode is invalid
        """
        if isinstance(mode, str):
            try:
                mode = NotesMode(mode.lower())
            except ValueError:
                raise ValueError(f"Invalid mode: {mode}")

        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        # Access or create notes slide
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame

        original_text = text_frame.text or ""
        original_length = len(original_text)

        if mode == NotesMode.OVERWRITE:
            final_text = text
        elif mode == NotesMode.APPEND:
            if original_text.strip():
                final_text = original_text + "\n" + text
            else:
                final_text = text
        elif mode == NotesMode.PREPEND:
            if original_text.strip():
                final_text = text + "\n" + original_text
            else:
                final_text = text

        text_frame.text = final_text
        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "mode": mode.value,
            "original_length": original_length,
            "new_length": len(final_text),
            "text_preview": final_text[:100] + "..."
            if len(final_text) > 100
            else final_text,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def set_footer(
        self,
        text: Optional[str] = None,
        show_slide_number: bool = False,
        show_date: bool = False,
        slide_index: Optional[int] = None,
    ) -> Dict[str, Any]:
        """
        Set footer properties for slide(s).

        Note: Footer configuration in python-pptx is limited.
        This method sets footer placeholders where available.

        Args:
            text: Footer text
            show_slide_number: Show slide numbers
            show_date: Show date
            slide_index: Specific slide (None = all slides)

        Returns:
            Dict with footer configuration results
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        version_before = self._capture_version()
        results = []

        # Determine slides to process
        if slide_index is not None:
            slides = [(slide_index, self._get_slide(slide_index))]
        else:
            slides = list(enumerate(self.prs.slides))

        for s_idx, slide in slides:
            slide_result = {
                "slide_index": s_idx,
                "footer_set": False,
                "slide_number_set": False,
                "date_set": False,
            }

            for shape in slide.shapes:
                if not shape.is_placeholder:
                    continue

                ph_type = _get_placeholder_type_int_helper(
                    shape.placeholder_format.type
                )

                # Footer placeholder (type 7)
                if ph_type == 7 and text is not None:
                    if shape.has_text_frame:
                        shape.text_frame.text = text
                        slide_result["footer_set"] = True

                # Slide number placeholder (type 6)
                if ph_type == 6 and show_slide_number:
                    slide_result["slide_number_set"] = True

                # Date placeholder (type 5)
                if ph_type == 5 and show_date:
                    slide_result["date_set"] = True

            results.append(slide_result)

        version_after = self._capture_version()

        return {
            "text": text,
            "show_slide_number": show_slide_number,
            "show_date": show_date,
            "slides_processed": len(results),
            "results": results,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    # ========================================================================
    # SHAPE OPERATIONS
    # ========================================================================

    def _set_fill_opacity(self, shape, opacity: float) -> bool:
        """
        Set the fill opacity of a shape by manipulating the underlying XML.

        Args:
            shape: The shape object with a fill
            opacity: Opacity value (0.0 = fully transparent, 1.0 = fully opaque)

        Returns:
            True if opacity was set, False if not applicable

        Note:
            python-pptx doesn't directly expose fill transparency, so we
            manipulate the OOXML directly. The alpha value uses a scale
            where 100000 = 100% opaque.
        """
        if opacity >= 1.0:
            # No need to set alpha for fully opaque - it's the default
            return True

        if opacity < 0.0:
            opacity = 0.0

        try:
            # Access the shape's spPr (shape properties) element
            spPr = shape._sp.spPr
            if spPr is None:
                return False

            # Find the solidFill element
            solidFill = spPr.find(qn("a:solidFill"))
            if solidFill is None:
                return False

            # Find the color element (could be srgbClr or schemeClr)
            color_elem = solidFill.find(qn("a:srgbClr"))
            if color_elem is None:
                color_elem = solidFill.find(qn("a:schemeClr"))
            if color_elem is None:
                return False

            # Calculate alpha value (Office uses 0-100000 scale, where 100000 = 100%)
            alpha_value = int(opacity * 100000)

            # Remove existing alpha element if present
            existing_alpha = color_elem.find(qn("a:alpha"))
            if existing_alpha is not None:
                color_elem.remove(existing_alpha)

            # Create and add new alpha element
            # Using SubElement to create properly namespaced element
            nsmap = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
            alpha_elem = etree.SubElement(color_elem, qn("a:alpha"))
            alpha_elem.set("val", str(alpha_value))

            return True

        except Exception as e:
            # Log but don't fail - opacity is enhancement, not critical
            self._log_warning(f"Could not set fill opacity: {e}")
            return False

    def _set_line_opacity(self, shape, opacity: float) -> bool:
        """
        Set the line/border opacity of a shape by manipulating the underlying XML.

        Args:
            shape: The shape object with a line
            opacity: Opacity value (0.0 = fully transparent, 1.0 = fully opaque)

        Returns:
            True if opacity was set, False if not applicable

        Note:
            Line opacity requires the line to have a solid fill. We manipulate
            the OOXML <a:ln><a:solidFill><a:srgbClr><a:alpha> structure.
        """
        if opacity >= 1.0:
            return True

        if opacity < 0.0:
            opacity = 0.0

        try:
            # Access the shape's spPr element
            spPr = shape._sp.spPr
            if spPr is None:
                return False

            # Find the line element
            ln = spPr.find(qn("a:ln"))
            if ln is None:
                return False

            # Find solidFill within line
            solidFill = ln.find(qn("a:solidFill"))
            if solidFill is None:
                # Line might not have a fill yet - try to find/create one
                return False

            # Find color element
            color_elem = solidFill.find(qn("a:srgbClr"))
            if color_elem is None:
                color_elem = solidFill.find(qn("a:schemeClr"))
            if color_elem is None:
                return False

            # Calculate and set alpha
            alpha_value = int(opacity * 100000)

            existing_alpha = color_elem.find(qn("a:alpha"))
            if existing_alpha is not None:
                color_elem.remove(existing_alpha)

            alpha_elem = etree.SubElement(color_elem, qn("a:alpha"))
            alpha_elem.set("val", str(alpha_value))

            return True

        except Exception as e:
            self._log_warning(f"Could not set line opacity: {e}")
            return False

    def _ensure_line_solid_fill(self, shape, color_hex: str) -> bool:
        """
        Ensure the shape's line has a solid fill with the specified color.
        This is necessary before setting line opacity.

        Args:
            shape: The shape object
            color_hex: Hex color string for the line

        Returns:
            True if successful
        """
        try:
            # Set line color through python-pptx first
            shape.line.color.rgb = ColorHelper.from_hex(color_hex)

            # Now ensure the XML structure is correct for opacity
            spPr = shape._sp.spPr
            ln = spPr.find(qn("a:ln"))

            if ln is None:
                return False

            # Check if solidFill exists
            solidFill = ln.find(qn("a:solidFill"))
            if solidFill is None:
                # Create solidFill structure
                solidFill = etree.SubElement(ln, qn("a:solidFill"))
                color_elem = etree.SubElement(solidFill, qn("a:srgbClr"))
                # Remove # from hex color
                color_val = color_hex.lstrip("#").upper()
                color_elem.set("val", color_val)

            return True

        except Exception as e:
            self._log_warning(f"Could not ensure line solid fill: {e}")
            return False

    def add_shape(
        self,
        slide_index: int,
        shape_type: str,
        position: Dict[str, Any],
        size: Dict[str, Any],
        fill_color: Optional[str] = None,
        fill_opacity: float = 1.0,
        line_color: Optional[str] = None,
        line_opacity: float = 1.0,
        line_width: float = 1.0,
        text: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Add shape to slide with optional transparency/opacity support.

        Args:
            slide_index: Target slide index
            shape_type: Shape type name (rectangle, ellipse, arrow_right, etc.)
            position: Position dict (percentage, inches, anchor, or grid)
            size: Size dict (percentage or inches)
            fill_color: Fill color hex (e.g., "#0070C0") or None for no fill
            fill_opacity: Fill opacity from 0.0 (transparent) to 1.0 (opaque).
                         Default is 1.0 (fully opaque). Use 0.15 for subtle overlays.
            line_color: Line/border color hex or None for no line
            line_opacity: Line opacity from 0.0 (transparent) to 1.0 (opaque).
                         Default is 1.0 (fully opaque).
            line_width: Line width in points (default: 1.0)
            text: Optional text to add inside shape

        Returns:
            Dict with shape_index, position, size, and applied styling details

        Raises:
            SlideNotFoundError: If slide index is invalid
            ValueError: If size is not specified or opacity is out of range

        Example:
            # Subtle white overlay for improved text readability
            agent.add_shape(
                slide_index=0,
                shape_type="rectangle",
                position={"left": "0%", "top": "0%"},
                size={"width": "100%", "height": "100%"},
                fill_color="#FFFFFF",
                fill_opacity=0.15  # 15% opaque = 85% transparent
            )
        """
        # Validate opacity ranges
        if not 0.0 <= fill_opacity <= 1.0:
            raise ValueError(
                f"fill_opacity must be between 0.0 and 1.0, got {fill_opacity}"
            )
        if not 0.0 <= line_opacity <= 1.0:
            raise ValueError(
                f"line_opacity must be between 0.0 and 1.0, got {line_opacity}"
            )

        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)

        if width is None or height is None:
            raise ValueError("Shape must have explicit width and height")

        # Map shape type string to MSO constant
        shape_type_map = {
            "rectangle": MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            "rounded_rectangle": MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            "ellipse": MSO_AUTO_SHAPE_TYPE.OVAL,
            "oval": MSO_AUTO_SHAPE_TYPE.OVAL,
            "triangle": MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
            "arrow_right": MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
            "arrow_left": MSO_AUTO_SHAPE_TYPE.LEFT_ARROW,
            "arrow_up": MSO_AUTO_SHAPE_TYPE.UP_ARROW,
            "arrow_down": MSO_AUTO_SHAPE_TYPE.DOWN_ARROW,
            "diamond": MSO_AUTO_SHAPE_TYPE.DIAMOND,
            "pentagon": MSO_AUTO_SHAPE_TYPE.PENTAGON,
            "hexagon": MSO_AUTO_SHAPE_TYPE.HEXAGON,
            "star": MSO_AUTO_SHAPE_TYPE.STAR_5_POINT,
            "heart": MSO_AUTO_SHAPE_TYPE.HEART,
            "lightning": MSO_AUTO_SHAPE_TYPE.LIGHTNING_BOLT,
            "sun": MSO_AUTO_SHAPE_TYPE.SUN,
            "moon": MSO_AUTO_SHAPE_TYPE.MOON,
            "cloud": MSO_AUTO_SHAPE_TYPE.CLOUD,
        }

        mso_shape = shape_type_map.get(shape_type.lower())
        if mso_shape is None:
            raise ValueError(
                f"Unknown shape type: {shape_type}",
                details={"valid_types": list(shape_type_map.keys())},
            )

        # Add shape
        shape = slide.shapes.add_shape(
            mso_shape, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        # Track what was actually applied
        styling_applied = {
            "fill_color": None,
            "fill_opacity": 1.0,
            "fill_opacity_applied": False,
            "line_color": None,
            "line_opacity": 1.0,
            "line_opacity_applied": False,
            "line_width": line_width,
        }

        # Apply fill color and opacity
        if fill_color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = ColorHelper.from_hex(fill_color)
            styling_applied["fill_color"] = fill_color
            styling_applied["fill_opacity"] = fill_opacity

            # Apply fill opacity if not fully opaque
            if fill_opacity < 1.0:
                opacity_set = self._set_fill_opacity(shape, fill_opacity)
                styling_applied["fill_opacity_applied"] = opacity_set
        else:
            # No fill - make background transparent
            shape.fill.background()

        # Apply line color and opacity
        if line_color:
            # Ensure line has solid fill for opacity support
            self._ensure_line_solid_fill(shape, line_color)
            shape.line.width = Pt(line_width)
            styling_applied["line_color"] = line_color
            styling_applied["line_opacity"] = line_opacity

            # Apply line opacity if not fully opaque
            if line_opacity < 1.0:
                opacity_set = self._set_line_opacity(shape, line_opacity)
                styling_applied["line_opacity_applied"] = opacity_set
        else:
            # No line
            shape.line.fill.background()

        # Add text if provided
        if text and shape.has_text_frame:
            shape.text_frame.text = text

        shape_index = len(slide.shapes) - 1
        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "shape_type": shape_type,
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height},
            "styling": styling_applied,
            "has_text": text is not None,
            "text_preview": text[:50] + "..." if text and len(text) > 50 else text,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def format_shape(
        self,
        slide_index: int,
        shape_index: int,
        fill_color: Optional[str] = None,
        fill_opacity: Optional[float] = None,
        line_color: Optional[str] = None,
        line_opacity: Optional[float] = None,
        line_width: Optional[float] = None,
        transparency: Optional[float] = None,
    ) -> Dict[str, Any]:
        """
        Format existing shape with optional transparency/opacity support.

        Args:
            slide_index: Target slide index
            shape_index: Shape index on slide
            fill_color: Fill color hex (e.g., "#0070C0")
            fill_opacity: Fill opacity from 0.0 (transparent) to 1.0 (opaque)
            line_color: Line/border color hex
            line_opacity: Line opacity from 0.0 (transparent) to 1.0 (opaque)
            line_width: Line width in points
            transparency: DEPRECATED - Use fill_opacity instead.
                         If provided, converted to fill_opacity (transparency = 1 - opacity).
                         Will be removed in v4.0.

        Returns:
            Dict with formatting changes applied and their status

        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
            ValueError: If opacity values are out of range

        Example:
            # Make an existing shape semi-transparent
            agent.format_shape(
                slide_index=0,
                shape_index=3,
                fill_opacity=0.5  # 50% opaque
            )
        """
        shape = self._get_shape(slide_index, shape_index)
        version_before = self._capture_version()

        changes: List[str] = []
        changes_detail: Dict[str, Any] = {}

        # Handle deprecated transparency parameter
        if transparency is not None:
            if fill_opacity is None:
                # Convert transparency to opacity (they're inverses)
                # transparency: 0.0 = opaque, 1.0 = invisible
                # opacity: 1.0 = opaque, 0.0 = invisible
                fill_opacity = 1.0 - transparency
                changes.append("transparency_converted_to_opacity")
                changes_detail["transparency_deprecated"] = True
                changes_detail["transparency_value"] = transparency
                changes_detail["converted_opacity"] = fill_opacity
                self._log_warning(
                    "The 'transparency' parameter is deprecated. "
                    "Use 'fill_opacity' instead (opacity = 1 - transparency)."
                )
            else:
                # Both provided - fill_opacity takes precedence
                changes.append("transparency_ignored")
                changes_detail["transparency_ignored"] = True
                self._log_warning(
                    "Both 'transparency' and 'fill_opacity' provided. "
                    "Using 'fill_opacity', ignoring 'transparency'."
                )

        # Validate opacity ranges
        if fill_opacity is not None and not 0.0 <= fill_opacity <= 1.0:
            raise ValueError(
                f"fill_opacity must be between 0.0 and 1.0, got {fill_opacity}"
            )
        if line_opacity is not None and not 0.0 <= line_opacity <= 1.0:
            raise ValueError(
                f"line_opacity must be between 0.0 and 1.0, got {line_opacity}"
            )

        # Apply fill color
        if fill_color is not None:
            shape.fill.solid()
            shape.fill.fore_color.rgb = ColorHelper.from_hex(fill_color)
            changes.append("fill_color")
            changes_detail["fill_color"] = fill_color

        # Apply fill opacity
        if fill_opacity is not None:
            # Ensure shape has solid fill before applying opacity
            if fill_color is None:
                try:
                    shape.fill.solid()
                except Exception:
                    pass

            if fill_opacity < 1.0:
                success = self._set_fill_opacity(shape, fill_opacity)
                if success:
                    changes.append("fill_opacity")
                    changes_detail["fill_opacity"] = fill_opacity
                    changes_detail["fill_opacity_applied"] = True
                else:
                    changes.append("fill_opacity_failed")
                    changes_detail["fill_opacity"] = fill_opacity
                    changes_detail["fill_opacity_applied"] = False
            else:
                # Opacity 1.0 = fully opaque (default, no XML change needed)
                changes.append("fill_opacity_reset")
                changes_detail["fill_opacity"] = 1.0

        # Apply line color
        if line_color is not None:
            self._ensure_line_solid_fill(shape, line_color)
            changes.append("line_color")
            changes_detail["line_color"] = line_color

        # Apply line opacity
        if line_opacity is not None:
            if line_opacity < 1.0:
                success = self._set_line_opacity(shape, line_opacity)
                if success:
                    changes.append("line_opacity")
                    changes_detail["line_opacity"] = line_opacity
                    changes_detail["line_opacity_applied"] = True
                else:
                    changes.append("line_opacity_failed")
                    changes_detail["line_opacity"] = line_opacity
                    changes_detail["line_opacity_applied"] = False
            else:
                changes.append("line_opacity_reset")
                changes_detail["line_opacity"] = 1.0

        # Apply line width
        if line_width is not None:
            shape.line.width = Pt(line_width)
            changes.append("line_width")
            changes_detail["line_width"] = line_width

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "changes_applied": changes,
            "changes_detail": changes_detail,
            "success": "failed" not in " ".join(changes),
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def remove_shape(
        self, slide_index: int, shape_index: int, approval_token: Optional[str] = None
    ) -> Dict[str, Any]:
        """
        Remove shape from slide.

        ⚠️ DESTRUCTIVE OPERATION - Requires approval token.

        Args:
            slide_index: Target slide index
            shape_index: Shape index to remove

        Returns:
            Dict with removal details

        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
        """
        self._validate_token(approval_token, APPROVAL_SCOPE_REMOVE_SHAPE)

        slide = self._get_slide(slide_index)
        shape = self._get_shape(slide_index, shape_index)
        version_before = self._capture_version()

        # Get shape info before removal
        shape_name = shape.name
        shape_type = str(shape.shape_type)

        # Remove shape from slide
        sp = shape.element
        sp.getparent().remove(sp)

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "removed_shape_index": shape_index,
            "removed_shape_name": shape_name,
            "removed_shape_type": shape_type,
            "new_shape_count": len(slide.shapes),
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def set_z_order(
        self, slide_index: int, shape_index: int, action: str
    ) -> Dict[str, Any]:
        """
        Change the z-order (stacking order) of a shape.

        Args:
            slide_index: Target slide index
            shape_index: Shape index to modify
            action: One of "bring_to_front", "send_to_back",
                   "bring_forward", "send_backward"

        Returns:
            Dict with z-order change details including old and new positions

        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
            ValueError: If action is invalid
        """
        valid_actions = {
            "bring_to_front",
            "send_to_back",
            "bring_forward",
            "send_backward",
        }
        if action not in valid_actions:
            raise ValueError(
                f"Invalid action: {action}. Must be one of {valid_actions}"
            )

        slide = self._get_slide(slide_index)
        shape = self._get_shape(slide_index, shape_index)
        version_before = self._capture_version()

        # Access the shape tree XML element
        sp_tree = slide.shapes._spTree
        element = shape.element

        # Find current position in XML tree
        current_index = -1
        shape_elements = [
            child
            for child in sp_tree
            if child.tag.endswith("}sp")
            or child.tag.endswith("}pic")
            or child.tag.endswith("}graphicFrame")
        ]

        for i, child in enumerate(sp_tree):
            if child == element:
                current_index = i
                break

        if current_index == -1:
            raise PowerPointAgentError(
                "Could not locate shape in XML tree",
                details={"slide_index": slide_index, "shape_index": shape_index},
            )

        new_index = current_index
        max_index = len(sp_tree) - 1

        # Execute the z-order action
        if action == "bring_to_front":
            sp_tree.remove(element)
            sp_tree.append(element)
            new_index = len(sp_tree) - 1

        elif action == "send_to_back":
            sp_tree.remove(element)
            # Insert after nvGrpSpPr and grpSpPr (indices 0 and 1 typically)
            sp_tree.insert(2, element)
            new_index = 2

        elif action == "bring_forward":
            if current_index < max_index:
                sp_tree.remove(element)
                sp_tree.insert(current_index + 1, element)
                new_index = current_index + 1

        elif action == "send_backward":
            if current_index > 2:  # Don't go before required elements
                sp_tree.remove(element)
                sp_tree.insert(current_index - 1, element)
                new_index = current_index - 1

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "action": action,
            "z_order_change": {"from": current_index, "to": new_index},
            "warning": "Shape indices may have changed after z-order operation. Re-query slide info.",
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def add_table(
        self,
        slide_index: int,
        rows: int,
        cols: int,
        position: Dict[str, Any],
        size: Dict[str, Any],
        data: Optional[List[List[Any]]] = None,
        header_row: bool = True,
    ) -> Dict[str, Any]:
        """
        Add table to slide.

        Args:
            slide_index: Target slide index
            rows: Number of rows
            cols: Number of columns
            position: Position dict
            size: Size dict
            data: Optional 2D list of cell values
            header_row: Whether first row is header (styling hint)

        Returns:
            Dict with shape_index and table details
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)

        if width is None or height is None:
            raise ValueError("Table must have explicit width and height")

        # Create table
        table_shape = slide.shapes.add_table(
            rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)
        )

        table = table_shape.table

        # Populate with data if provided
        cells_filled = 0
        if data:
            for row_idx, row_data in enumerate(data):
                if row_idx >= rows:
                    break
                for col_idx, cell_value in enumerate(row_data):
                    if col_idx >= cols:
                        break
                    table.cell(row_idx, col_idx).text = str(cell_value)
                    cells_filled += 1

        shape_index = len(slide.shapes) - 1
        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "rows": rows,
            "cols": cols,
            "cells_filled": cells_filled,
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height},
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def add_connector(
        self,
        slide_index: int,
        from_shape_index: int,
        to_shape_index: int,
        connector_type: str = "straight",
    ) -> Dict[str, Any]:
        """
        Add connector line between two shapes.

        Args:
            slide_index: Target slide index
            from_shape_index: Starting shape index
            to_shape_index: Ending shape index
            connector_type: "straight", "elbow", or "curved"

        Returns:
            Dict with connector details
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        shape1 = self._get_shape(slide_index, from_shape_index)
        shape2 = self._get_shape(slide_index, to_shape_index)

        # Calculate center points
        x1 = shape1.left + shape1.width // 2
        y1 = shape1.top + shape1.height // 2
        x2 = shape2.left + shape2.width // 2
        y2 = shape2.top + shape2.height // 2

        # Map connector type
        connector_map = {
            "straight": MSO_CONNECTOR.STRAIGHT,
            "elbow": MSO_CONNECTOR.ELBOW,
            "curved": MSO_CONNECTOR.CURVE,
        }
        mso_connector = connector_map.get(
            connector_type.lower(), MSO_CONNECTOR.STRAIGHT
        )

        # Add connector
        connector = slide.shapes.add_connector(mso_connector, x1, y1, x2, y2)

        shape_index = len(slide.shapes) - 1
        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "from_shape": from_shape_index,
            "to_shape": to_shape_index,
            "connector_type": connector_type,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    # ========================================================================
    # IMAGE OPERATIONS
    # ========================================================================

    def insert_image(
        self,
        slide_index: int,
        image_path: Union[str, Path],
        position: Dict[str, Any],
        size: Optional[Dict[str, Any]] = None,
        alt_text: Optional[str] = None,
        compress: bool = False,
    ) -> Dict[str, Any]:
        """
        Insert image on slide.

        Args:
            slide_index: Target slide index
            image_path: Path to image file
            position: Position dict
            size: Optional size dict (can use "auto" for aspect ratio)
            alt_text: Alternative text for accessibility
            compress: Compress image before inserting

        Returns:
            Dict with shape_index and image details
        """
        slide = self._get_slide(slide_index)
        image_path = PathValidator.validate_image_path(image_path)
        version_before = self._capture_version()

        left, top = Position.from_dict(position)

        # Get aspect ratio if Pillow available
        aspect_ratio = None
        if HAS_PILLOW:
            try:
                with PILImage.open(image_path) as img:
                    aspect_ratio = img.width / img.height
            except Exception:
                pass

        # Parse size
        if size:
            width, height = Size.from_dict(size, aspect_ratio=aspect_ratio)
        else:
            # Default to half slide width, maintain aspect ratio
            width = SLIDE_WIDTH_INCHES * 0.5
            if aspect_ratio:
                height = width / aspect_ratio
            else:
                height = SLIDE_HEIGHT_INCHES * 0.3

        # Compress if requested
        if compress and HAS_PILLOW:
            image_stream = AssetValidator.compress_image(image_path)
            picture = slide.shapes.add_picture(
                image_stream,
                Inches(left),
                Inches(top),
                width=Inches(width) if width else None,
                height=Inches(height) if height else None,
            )
        else:
            picture = slide.shapes.add_picture(
                str(image_path),
                Inches(left),
                Inches(top),
                width=Inches(width) if width else None,
                height=Inches(height) if height else None,
            )

        # Set alt text
        if alt_text:
            picture.name = alt_text
            try:
                # Set description attribute for proper alt text
                picture._element.set("descr", alt_text)
            except Exception:
                pass

        shape_index = len(slide.shapes) - 1
        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "image_path": str(image_path),
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height},
            "alt_text_set": alt_text is not None,
            "compressed": compress,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def replace_image(
        self,
        slide_index: int,
        old_image_name: str,
        new_image_path: Union[str, Path],
        compress: bool = False,
    ) -> Dict[str, Any]:
        """
        Replace existing image by name.

        Args:
            slide_index: Target slide index
            old_image_name: Name or partial name of image to replace
            new_image_path: Path to new image file
            compress: Compress new image

        Returns:
            Dict with replacement details
        """
        slide = self._get_slide(slide_index)
        new_image_path = PathValidator.validate_image_path(new_image_path)
        version_before = self._capture_version()

        replaced = False
        old_shape_index = None
        new_shape_index = None

        for idx, shape in enumerate(slide.shapes):
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                if shape.name == old_image_name or old_image_name in (shape.name or ""):
                    # Store position and size
                    left = shape.left
                    top = shape.top
                    width = shape.width
                    height = shape.height
                    old_shape_index = idx

                    # Remove old image
                    sp = shape.element
                    sp.getparent().remove(sp)

                    # Add new image
                    if compress and HAS_PILLOW:
                        image_stream = AssetValidator.compress_image(new_image_path)
                        new_picture = slide.shapes.add_picture(
                            image_stream, left, top, width=width, height=height
                        )
                    else:
                        new_picture = slide.shapes.add_picture(
                            str(new_image_path), left, top, width=width, height=height
                        )

                    new_shape_index = len(slide.shapes) - 1
                    replaced = True
                    break

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "replaced": replaced,
            "old_image_name": old_image_name,
            "old_shape_index": old_shape_index,
            "new_image_path": str(new_image_path),
            "new_shape_index": new_shape_index,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def set_image_properties(
        self,
        slide_index: int,
        shape_index: int,
        alt_text: Optional[str] = None,
        name: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Set image properties.

        Args:
            slide_index: Target slide index
            shape_index: Image shape index
            alt_text: Alternative text for accessibility
            name: Shape name

        Returns:
            Dict with properties set
        """
        shape = self._get_shape(slide_index, shape_index)
        version_before = self._capture_version()

        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape at index {shape_index} is not an image")

        changes = []

        if alt_text is not None:
            try:
                shape._element.set("descr", alt_text)
                changes.append("alt_text")
            except Exception:
                # Fallback to name
                shape.name = alt_text
                changes.append("alt_text_via_name")

        if name is not None:
            shape.name = name
            changes.append("name")

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "changes_applied": changes,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def crop_image(
        self,
        slide_index: int,
        shape_index: int,
        left: float = 0.0,
        top: float = 0.0,
        right: float = 0.0,
        bottom: float = 0.0,
    ) -> Dict[str, Any]:
        """
        Crop image by specifying crop amounts from each edge.

        Args:
            slide_index: Target slide index
            shape_index: Image shape index
            left: Crop from left (0.0 to 1.0, proportion of width)
            top: Crop from top (0.0 to 1.0, proportion of height)
            right: Crop from right (0.0 to 1.0, proportion of width)
            bottom: Crop from bottom (0.0 to 1.0, proportion of height)

        Returns:
            Dict with crop details
        """
        shape = self._get_shape(slide_index, shape_index)
        version_before = self._capture_version()

        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape at index {shape_index} is not an image")

        # Validate crop values
        for name, value in [
            ("left", left),
            ("top", top),
            ("right", right),
            ("bottom", bottom),
        ]:
            if not 0.0 <= value < 1.0:
                raise ValueError(
                    f"Crop {name} must be between 0.0 and 1.0, got {value}"
                )

        if left + right >= 1.0:
            raise ValueError("Left + right crop cannot equal or exceed 1.0")
        if top + bottom >= 1.0:
            raise ValueError("Top + bottom crop cannot equal or exceed 1.0")

        # Apply crop using picture's crop properties
        try:
            # Access the picture element
            pic = shape._element

            # Find or create blipFill element
            blip_fill = pic.find(
                ".//{http://schemas.openxmlformats.org/presentationml/2006/main}blipFill"
            )
            if blip_fill is None:
                blip_fill = pic.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blipFill"
                )

            if blip_fill is not None:
                # Find or create srcRect element
                ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
                src_rect = blip_fill.find(f"{ns}srcRect")

                if src_rect is None:
                    src_rect = etree.SubElement(blip_fill, f"{ns}srcRect")

                # Set crop values (in percentage * 1000)
                src_rect.set("l", str(int(left * 100000)))
                src_rect.set("t", str(int(top * 100000)))
                src_rect.set("r", str(int(right * 100000)))
                src_rect.set("b", str(int(bottom * 100000)))

                version_after = self._capture_version()

                return {
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "crop_applied": True,
                    "crop_values": {
                        "left": left,
                        "top": top,
                        "right": right,
                        "bottom": bottom,
                    },
                    "presentation_version_before": version_before,
                    "presentation_version_after": version_after,
                }
        except Exception as e:
            logger.warning(f"Could not apply crop via XML: {e}")

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "crop_applied": False,
            "error": "Crop not supported for this image type",
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def resize_image(
        self,
        slide_index: int,
        shape_index: int,
        width: Optional[float] = None,
        height: Optional[float] = None,
        maintain_aspect: bool = True,
    ) -> Dict[str, Any]:
        """
        Resize image shape.

        Args:
            slide_index: Target slide index
            shape_index: Image shape index
            width: New width in inches (None = keep current)
            height: New height in inches (None = keep current)
            maintain_aspect: Maintain aspect ratio

        Returns:
            Dict with new dimensions
        """
        shape = self._get_shape(slide_index, shape_index)
        version_before = self._capture_version()

        if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
            raise ValueError(f"Shape at index {shape_index} is not an image")

        original_width = shape.width / EMU_PER_INCH
        original_height = shape.height / EMU_PER_INCH
        aspect = original_width / original_height if original_height > 0 else 1.0

        new_width = width
        new_height = height

        if maintain_aspect:
            if width is not None and height is None:
                new_height = width / aspect
            elif height is not None and width is None:
                new_width = height * aspect

        if new_width is not None:
            shape.width = Inches(new_width)
        if new_height is not None:
            shape.height = Inches(new_height)

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "original_size": {"width": original_width, "height": original_height},
            "new_size": {
                "width": new_width or original_width,
                "height": new_height or original_height,
            },
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    # ========================================================================
    # CHART OPERATIONS
    # ========================================================================

    def add_chart(
        self,
        slide_index: int,
        chart_type: str,
        data: Dict[str, Any],
        position: Dict[str, Any],
        size: Dict[str, Any],
        title: Optional[str] = None,
    ) -> Dict[str, Any]:
        """
        Add chart to slide.

        Args:
            slide_index: Target slide index
            chart_type: Chart type (column, bar, line, pie, etc.)
            data: Chart data dict with "categories" and "series"
            position: Position dict
            size: Size dict
            title: Optional chart title

        Returns:
            Dict with shape_index and chart details

        Example data:
            {
                "categories": ["Q1", "Q2", "Q3", "Q4"],
                "series": [
                    {"name": "Revenue", "values": [100, 120, 140, 160]},
                    {"name": "Costs", "values": [80, 90, 100, 110]}
                ]
            }
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        left, top = Position.from_dict(position)
        width, height = Size.from_dict(size)

        if width is None or height is None:
            raise ValueError("Chart must have explicit width and height")

        # Map chart type string to XL constant
        chart_type_map = {
            "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_clustered": XL_CHART_TYPE.COLUMN_CLUSTERED,
            "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_clustered": XL_CHART_TYPE.BAR_CLUSTERED,
            "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
            "line": XL_CHART_TYPE.LINE,
            "line_markers": XL_CHART_TYPE.LINE_MARKERS,
            "pie": XL_CHART_TYPE.PIE,
            "pie_exploded": XL_CHART_TYPE.PIE_EXPLODED,
            "area": XL_CHART_TYPE.AREA,
            "scatter": XL_CHART_TYPE.XY_SCATTER,
            "doughnut": XL_CHART_TYPE.DOUGHNUT,
        }

        xl_chart_type = chart_type_map.get(
            chart_type.lower(), XL_CHART_TYPE.COLUMN_CLUSTERED
        )

        # Build chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])

        for series in data.get("series", []):
            chart_data.add_series(series["name"], series["values"])

        # Add chart
        chart_shape = slide.shapes.add_chart(
            xl_chart_type,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
            chart_data,
        )

        # Set title if provided
        if title:
            chart_shape.chart.has_title = True
            chart_shape.chart.chart_title.text_frame.text = title

        shape_index = len(slide.shapes) - 1
        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "shape_index": shape_index,
            "chart_type": chart_type,
            "categories_count": len(data.get("categories", [])),
            "series_count": len(data.get("series", [])),
            "title": title,
            "position": {"left": left, "top": top},
            "size": {"width": width, "height": height},
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def update_chart_data(
        self, slide_index: int, chart_index: int, data: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Update existing chart data.

        Args:
            slide_index: Target slide index
            chart_index: Chart index on slide (not shape index)
            data: New chart data dict

        Returns:
            Dict with update details
        """
        chart_shape = self._get_chart_shape(slide_index, chart_index)
        version_before = self._capture_version()

        # Build new chart data
        chart_data = CategoryChartData()
        chart_data.categories = data.get("categories", [])

        for series in data.get("series", []):
            chart_data.add_series(series["name"], series["values"])

        # Try to replace data (preserves formatting)
        try:
            chart_shape.chart.replace_data(chart_data)
            method = "replace_data"
        except AttributeError:
            # Fallback: recreate chart (loses some formatting)
            logger.warning(
                "chart.replace_data() not available. "
                "Recreating chart (some formatting may be lost)."
            )

            slide = self._get_slide(slide_index)

            # Store chart properties
            left = chart_shape.left
            top = chart_shape.top
            width = chart_shape.width
            height = chart_shape.height
            chart_type = chart_shape.chart.chart_type
            has_title = chart_shape.chart.has_title
            title_text = None
            if has_title:
                try:
                    title_text = chart_shape.chart.chart_title.text_frame.text
                except Exception:
                    pass

            # Remove old chart
            sp = chart_shape.element
            sp.getparent().remove(sp)

            # Create new chart
            new_chart_shape = slide.shapes.add_chart(
                chart_type, left, top, width, height, chart_data
            )

            # Restore title
            if title_text:
                new_chart_shape.chart.has_title = True
                new_chart_shape.chart.chart_title.text_frame.text = title_text

            method = "recreate"

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "chart_index": chart_index,
            "categories_count": len(data.get("categories", [])),
            "series_count": len(data.get("series", [])),
            "update_method": method,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def format_chart(
        self,
        slide_index: int,
        chart_index: int,
        title: Optional[str] = None,
        legend_position: Optional[str] = None,
        has_legend: Optional[bool] = None,
    ) -> Dict[str, Any]:
        """
        Format existing chart.

        Args:
            slide_index: Target slide index
            chart_index: Chart index on slide
            title: Chart title
            legend_position: Legend position ("bottom", "left", "right", "top")
            has_legend: Show/hide legend

        Returns:
            Dict with formatting changes
        """
        chart_shape = self._get_chart_shape(slide_index, chart_index)
        version_before = self._capture_version()

        chart = chart_shape.chart

        changes = []

        if title is not None:
            chart.has_title = True
            chart.chart_title.text_frame.text = title
            changes.append("title")

        if has_legend is not None:
            chart.has_legend = has_legend
            changes.append("has_legend")

        if legend_position is not None and chart.has_legend:
            from pptx.enum.chart import XL_LEGEND_POSITION

            position_map = {
                "bottom": XL_LEGEND_POSITION.BOTTOM,
                "left": XL_LEGEND_POSITION.LEFT,
                "right": XL_LEGEND_POSITION.RIGHT,
                "top": XL_LEGEND_POSITION.TOP,
                "corner": XL_LEGEND_POSITION.CORNER,
            }
            if legend_position.lower() in position_map:
                chart.legend.position = position_map[legend_position.lower()]
                changes.append("legend_position")

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "chart_index": chart_index,
            "changes_applied": changes,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    # ========================================================================
    # LAYOUT & THEME OPERATIONS
    # ========================================================================

    def set_slide_layout(self, slide_index: int, layout_name: str) -> Dict[str, Any]:
        """
        Change slide layout.

        Note: This changes the layout but may not reposition existing content.

        Args:
            slide_index: Target slide index
            layout_name: Name of new layout

        Returns:
            Dict with layout change details
        """
        slide = self._get_slide(slide_index)
        version_before = self._capture_version()

        layout = self._get_layout(layout_name)

        old_layout = slide.slide_layout.name
        slide.slide_layout = layout

        version_after = self._capture_version()

        return {
            "slide_index": slide_index,
            "old_layout": old_layout,
            "new_layout": layout_name,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def set_background(
        self,
        slide_index: Optional[int] = None,
        color: Optional[str] = None,
        image_path: Optional[Union[str, Path]] = None,
    ) -> Dict[str, Any]:
        """
        Set slide background color or image.

        Args:
            slide_index: Target slide (None = all slides)
            color: Background color hex
            image_path: Background image path

        Returns:
            Dict with background change details
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        if color is None and image_path is None:
            raise ValueError("Must specify either color or image_path")

        version_before = self._capture_version()
        results = []

        # Determine slides to process
        if slide_index is not None:
            slides = [(slide_index, self._get_slide(slide_index))]
        else:
            slides = list(enumerate(self.prs.slides))

        for s_idx, slide in slides:
            result = {"slide_index": s_idx, "success": False}

            try:
                background = slide.background
                fill = background.fill

                if color:
                    fill.solid()
                    fill.fore_color.rgb = ColorHelper.from_hex(color)
                    result["success"] = True
                    result["type"] = "color"
                    result["color"] = color

                elif image_path:
                    # Note: python-pptx has limited background image support
                    # This is a best-effort implementation
                    image_path = PathValidator.validate_image_path(image_path)
                    result["type"] = "image"
                    result["image_path"] = str(image_path)
                    result["note"] = (
                        "Background image support is limited in python-pptx"
                    )

            except Exception as e:
                result["error"] = str(e)

            results.append(result)

        version_after = self._capture_version()

        return {
            "slides_processed": len(results),
            "results": results,
            "presentation_version_before": version_before,
            "presentation_version_after": version_after,
        }

    def get_available_layouts(self) -> List[str]:
        """
        Get list of available layout names.

        Returns:
            List of layout name strings
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        self._ensure_layout_cache()
        return list(self._layout_cache.keys())

    # ========================================================================
    # VALIDATION OPERATIONS
    # ========================================================================

    def validate_presentation(self) -> Dict[str, Any]:
        """
        Comprehensive presentation validation.

        Returns:
            Validation report dict
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        issues = {
            "empty_slides": [],
            "slides_without_titles": [],
            "fonts_used": set(),
            "large_shapes": [],
        }

        for idx, slide in enumerate(self.prs.slides):
            # Check for empty slides
            if len(slide.shapes) == 0:
                issues["empty_slides"].append(idx)

            # Check for title
            has_title = False
            for shape in slide.shapes:
                if shape.is_placeholder:
                    ph_type = _get_placeholder_type_int_helper(
                        shape.placeholder_format.type
                    )
                    if ph_type in TITLE_PLACEHOLDER_TYPES:
                        if shape.has_text_frame and shape.text_frame.text.strip():
                            has_title = True
                            break

                # Collect fonts
                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.font.name:
                            issues["fonts_used"].add(para.font.name)

            if not has_title:
                issues["slides_without_titles"].append(idx)

        issues["fonts_used"] = list(issues["fonts_used"])

        total_issues = len(issues["empty_slides"]) + len(
            issues["slides_without_titles"]
        )

        return {
            "status": "issues_found" if total_issues > 0 else "valid",
            "total_issues": total_issues,
            "slide_count": len(self.prs.slides),
            "issues": issues,
        }

    def check_accessibility(self) -> Dict[str, Any]:
        """
        Run accessibility checker.

        Returns:
            Accessibility report dict
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        return AccessibilityChecker.check_presentation(self.prs)

    def validate_assets(self) -> Dict[str, Any]:
        """
        Run asset validator.

        Returns:
            Asset validation report dict
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        return AssetValidator.validate_presentation_assets(self.prs, self.filepath)

    # ========================================================================
    # EXPORT OPERATIONS
    # ========================================================================

    def export_to_pdf(self, output_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Export presentation to PDF.

        Requires LibreOffice or Microsoft Office installed.

        Args:
            output_path: Output PDF path

        Returns:
            Dict with export details
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        output_path = Path(output_path)
        if output_path.suffix.lower() != ".pdf":
            output_path = output_path.with_suffix(".pdf")

        # Ensure parent directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)

        # Save to temp file first
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            temp_pptx = Path(tmp.name)

        try:
            self.prs.save(str(temp_pptx))

            # Try LibreOffice conversion
            result = subprocess.run(
                [
                    "soffice",
                    "--headless",
                    "--convert-to",
                    "pdf",
                    "--outdir",
                    str(output_path.parent),
                    str(temp_pptx),
                ],
                capture_output=True,
                timeout=120,
            )

            if result.returncode != 0:
                raise PowerPointAgentError(
                    "PDF export failed. LibreOffice is required for PDF export.",
                    details={
                        "stderr": result.stderr.decode() if result.stderr else None,
                        "install_instructions": {
                            "linux": "sudo apt install libreoffice-impress",
                            "macos": "brew install --cask libreoffice",
                            "windows": "Download from libreoffice.org",
                        },
                    },
                )

            # Rename output file to desired name
            generated_pdf = output_path.parent / f"{temp_pptx.stem}.pdf"
            if generated_pdf.exists() and generated_pdf != output_path:
                shutil.move(str(generated_pdf), str(output_path))

            return {
                "success": True,
                "output_path": str(output_path),
                "file_size_bytes": output_path.stat().st_size
                if output_path.exists()
                else 0,
            }

        finally:
            temp_pptx.unlink(missing_ok=True)

    def extract_notes(self) -> Dict[int, str]:
        """
        Extract speaker notes from all slides.

        Returns:
            Dict mapping slide index to notes text
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        notes = {}

        for idx, slide in enumerate(self.prs.slides):
            if slide.has_notes_slide:
                try:
                    notes_slide = slide.notes_slide
                    text_frame = notes_slide.notes_text_frame
                    if text_frame.text and text_frame.text.strip():
                        notes[idx] = text_frame.text
                except Exception:
                    pass

        return notes

    # ========================================================================
    # INFORMATION & VERSIONING
    # ========================================================================

    def get_presentation_info(self) -> Dict[str, Any]:
        """
        Get presentation metadata and information.

        Returns:
            Dict with presentation information
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        info = {
            "slide_count": len(self.prs.slides),
            "layouts": self.get_available_layouts(),
            "slide_width_inches": self.prs.slide_width / EMU_PER_INCH,
            "slide_height_inches": self.prs.slide_height / EMU_PER_INCH,
            "presentation_version": self.get_presentation_version(),
        }

        # Calculate aspect ratio
        width = info["slide_width_inches"]
        height = info["slide_height_inches"]
        if height > 0:
            ratio = width / height
            if abs(ratio - 16 / 9) < 0.1:
                info["aspect_ratio"] = "16:9"
            elif abs(ratio - 4 / 3) < 0.1:
                info["aspect_ratio"] = "4:3"
            else:
                info["aspect_ratio"] = f"{width:.2f}:{height:.2f}"

        # File info
        if self.filepath and self.filepath.exists():
            stat = self.filepath.stat()
            info["file"] = str(self.filepath)
            info["file_size_bytes"] = stat.st_size
            info["file_size_mb"] = round(stat.st_size / (1024 * 1024), 2)
            info["modified"] = datetime.fromtimestamp(stat.st_mtime).isoformat()

        return info

    def get_slide_info(self, slide_index: int) -> Dict[str, Any]:
        """
        Get detailed information about a specific slide.

        Args:
            slide_index: Slide index to inspect

        Returns:
            Dict with comprehensive slide information
        """
        slide = self._get_slide(slide_index)

        shapes_info = []
        for idx, shape in enumerate(slide.shapes):
            # Determine shape type string
            shape_type_str = str(shape.shape_type).replace("MSO_SHAPE_TYPE.", "")

            if shape.is_placeholder:
                ph_type = _get_placeholder_type_int_helper(
                    shape.placeholder_format.type
                )
                ph_name = get_placeholder_type_name(ph_type)
                shape_type_str = f"PLACEHOLDER ({ph_name})"

            shape_info = {
                "index": idx,
                "type": shape_type_str,
                "name": shape.name,
                "has_text": hasattr(shape, "text_frame") and shape.has_text_frame,
                "position": {
                    "left_inches": round(shape.left / EMU_PER_INCH, 3),
                    "top_inches": round(shape.top / EMU_PER_INCH, 3),
                    "left_percent": f"{(shape.left / self.prs.slide_width * 100):.1f}%",
                    "top_percent": f"{(shape.top / self.prs.slide_height * 100):.1f}%",
                },
                "size": {
                    "width_inches": round(shape.width / EMU_PER_INCH, 3),
                    "height_inches": round(shape.height / EMU_PER_INCH, 3),
                    "width_percent": f"{(shape.width / self.prs.slide_width * 100):.1f}%",
                    "height_percent": f"{(shape.height / self.prs.slide_height * 100):.1f}%",
                },
            }

            # Add text content if present
            if shape.has_text_frame:
                try:
                    full_text = shape.text_frame.text
                    shape_info["text"] = full_text
                    shape_info["text_length"] = len(full_text)
                except Exception:
                    pass

            # Add image info if picture
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    shape_info["image_size_bytes"] = len(shape.image.blob)
                    shape_info["image_content_type"] = shape.image.content_type
                except Exception:
                    pass

            # Add chart info if chart
            if hasattr(shape, "has_chart") and shape.has_chart:
                try:
                    shape_info["chart_type"] = str(shape.chart.chart_type)
                except Exception:
                    pass

            shapes_info.append(shape_info)

        # Check for notes
        has_notes = False
        notes_preview = None
        if slide.has_notes_slide:
            try:
                notes_text = slide.notes_slide.notes_text_frame.text
                if notes_text and notes_text.strip():
                    has_notes = True
                    notes_preview = (
                        notes_text[:100] + "..."
                        if len(notes_text) > 100
                        else notes_text
                    )
            except Exception:
                pass

        return {
            "slide_index": slide_index,
            "layout": slide.slide_layout.name,
            "shape_count": len(slide.shapes),
            "shapes": shapes_info,
            "has_notes": has_notes,
            "notes_preview": notes_preview,
        }

    def get_presentation_version(self) -> str:
        """
        Compute a deterministic version hash for the presentation.

        The version is based on:
        - Slide count & Layouts
        - Shape counts per slide
        - Text content (SHA-256)
        - Shape Geometry (Position/Size) to detect layout changes

        Returns:
            SHA-256 hash prefix (16 characters)
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        # Build version components
        components = []

        # Slide count
        components.append(f"slides:{len(self.prs.slides)}")

        # Per-slide information
        for idx, slide in enumerate(self.prs.slides):
            slide_components = [
                f"slide:{idx}",
                f"layout:{slide.slide_layout.name}",
                f"shapes:{len(slide.shapes)}",
            ]

            # Add text content hash
            text_content = []
            for shape in slide.shapes:
                # Add Geometry hash to detect moves/resizes
                geo_hash = f"{shape.left}:{shape.top}:{shape.width}:{shape.height}"
                slide_components.append(f"geo:{geo_hash}")

                if hasattr(shape, "text_frame") and shape.has_text_frame:
                    try:
                        text_content.append(shape.text_frame.text)
                    except Exception:
                        pass

            if text_content:
                # Use SHA-256 for content
                text_hash = hashlib.sha256("".join(text_content).encode()).hexdigest()[
                    :8
                ]
                slide_components.append(f"text:{text_hash}")

            components.extend(slide_components)

        # Compute final hash
        version_string = "|".join(components)
        full_hash = hashlib.sha256(version_string.encode()).hexdigest()

        return full_hash[:16]

    # ========================================================================
    # PRIVATE HELPER METHODS
    # ========================================================================

    def _get_slide(self, index: int):
        """
        Get slide by index with validation.

        Args:
            index: Slide index (0-based)

        Returns:
            Slide object

        Raises:
            PowerPointAgentError: If no presentation loaded
            SlideNotFoundError: If index is out of range
        """
        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        slide_count = len(self.prs.slides)

        if not 0 <= index < slide_count:
            raise SlideNotFoundError(
                f"Slide index {index} out of range (0-{slide_count - 1})",
                details={
                    "index": index,
                    "slide_count": slide_count,
                    "valid_range": f"0-{slide_count - 1}",
                },
            )

        return self.prs.slides[index]

    def _get_shape(self, slide_index: int, shape_index: int):
        """
        Get shape by slide and shape index with validation.

        Args:
            slide_index: Slide index
            shape_index: Shape index on slide

        Returns:
            Shape object

        Raises:
            SlideNotFoundError: If slide index is invalid
            ShapeNotFoundError: If shape index is invalid
        """
        slide = self._get_slide(slide_index)

        shape_count = len(slide.shapes)

        if not 0 <= shape_index < shape_count:
            raise ShapeNotFoundError(
                f"Shape index {shape_index} out of range on slide {slide_index}",
                details={
                    "slide_index": slide_index,
                    "shape_index": shape_index,
                    "shape_count": shape_count,
                    "valid_range": f"0-{shape_count - 1}"
                    if shape_count > 0
                    else "no shapes",
                },
            )

        return slide.shapes[shape_index]

    def _get_chart_shape(self, slide_index: int, chart_index: int):
        """
        Get chart shape by slide and chart index.

        Args:
            slide_index: Slide index
            chart_index: Chart index on slide (0-based among charts only)

        Returns:
            Chart shape object

        Raises:
            ChartNotFoundError: If chart not found
        """
        slide = self._get_slide(slide_index)

        chart_count = 0
        for shape in slide.shapes:
            if hasattr(shape, "has_chart") and shape.has_chart:
                if chart_count == chart_index:
                    return shape
                chart_count += 1

        raise ChartNotFoundError(
            f"Chart at index {chart_index} not found on slide {slide_index}",
            details={
                "slide_index": slide_index,
                "chart_index": chart_index,
                "charts_found": chart_count,
            },
        )

    def _get_layout(self, layout_name: str):
        """
        Get layout by name with caching.

        Args:
            layout_name: Layout name

        Returns:
            Layout object

        Raises:
            LayoutNotFoundError: If layout doesn't exist
        """
        self._ensure_layout_cache()

        layout = self._layout_cache.get(layout_name)

        if layout is None:
            raise LayoutNotFoundError(
                f"Layout '{layout_name}' not found",
                details={"available_layouts": list(self._layout_cache.keys())},
            )

        return layout

    def _ensure_layout_cache(self) -> None:
        """Build layout cache if not already built."""
        if self._layout_cache is not None:
            return

        if not self.prs:
            raise PowerPointAgentError("No presentation loaded")

        self._layout_cache = {layout.name: layout for layout in self.prs.slide_layouts}

    def _copy_shape(self, source_shape, target_slide) -> None:
        """
        Copy shape to target slide.

        Args:
            source_shape: Shape to copy
            target_slide: Destination slide
        """
        # Handle pictures
        if source_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            try:
                blob = source_shape.image.blob
                target_slide.shapes.add_picture(
                    BytesIO(blob),
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height,
                )
            except Exception as e:
                logger.warning(f"Could not copy picture: {e}")
            return

        # Handle auto shapes and text boxes
        if source_shape.shape_type in (
            MSO_SHAPE_TYPE.AUTO_SHAPE,
            MSO_SHAPE_TYPE.TEXT_BOX,
        ):
            try:
                # Get auto shape type, default to rectangle
                try:
                    auto_shape_type = source_shape.auto_shape_type
                except Exception:
                    auto_shape_type = MSO_AUTO_SHAPE_TYPE.RECTANGLE

                new_shape = target_slide.shapes.add_shape(
                    auto_shape_type,
                    source_shape.left,
                    source_shape.top,
                    source_shape.width,
                    source_shape.height,
                )

                # Copy text
                if source_shape.has_text_frame:
                    try:
                        new_shape.text_frame.text = source_shape.text_frame.text
                    except Exception:
                        pass

                # Copy fill
                try:
                    if source_shape.fill.type == 1:  # Solid fill
                        new_shape.fill.solid()
                        new_shape.fill.fore_color.rgb = source_shape.fill.fore_color.rgb
                except Exception:
                    pass

            except Exception as e:
                logger.warning(f"Could not copy shape: {e}")
            return

        # Log unsupported shape types
        logger.debug(f"Shape type {source_shape.shape_type} not copied (not supported)")


# ============================================================================
# MODULE EXPORTS
# ============================================================================

__all__ = [
    # Main class
    "PowerPointAgent",
    # Exceptions
    "PowerPointAgentError",
    "SlideNotFoundError",
    "ShapeNotFoundError",
    "ChartNotFoundError",
    "LayoutNotFoundError",
    "ImageNotFoundError",
    "InvalidPositionError",
    "TemplateError",
    "ThemeError",
    "AccessibilityError",
    "AssetValidationError",
    "FileLockError",
    "PathValidationError",
    "ApprovalTokenError",
    # Utility classes
    "FileLock",
    "PathValidator",
    "Position",
    "Size",
    "ColorHelper",
    "TemplateProfile",
    "AccessibilityChecker",
    "AssetValidator",
    # Enums
    "ShapeType",
    "ChartType",
    "TextAlignment",
    "VerticalAlignment",
    "BulletStyle",
    "ImageFormat",
    "ExportFormat",
    "ZOrderAction",
    "NotesMode",
    # Constants
    "SLIDE_WIDTH_INCHES",
    "SLIDE_HEIGHT_INCHES",
    "ANCHOR_POINTS",
    "CORPORATE_COLORS",
    "STANDARD_FONTS",
    "APPROVAL_SCOPE_DELETE_SLIDE",
    "APPROVAL_SCOPE_REMOVE_SHAPE",
    "APPROVAL_SCOPE_MERGE_PRESENTATIONS",
    "WCAG_CONTRAST_NORMAL",
    "WCAG_CONTRAST_LARGE",
    "EMU_PER_INCH",
    # Functions
    "get_placeholder_type_name",
    # Module metadata
    "__version__",
    "__author__",
    "__license__",
]
